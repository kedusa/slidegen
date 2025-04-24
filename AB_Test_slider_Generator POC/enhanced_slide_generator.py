import io
import tempfile
import os
import re
import random
import numpy as np
import cv2
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pytesseract

class EnhancedSlideGenerator:
    """
    Enhanced slide generator that creates professional A/B test slides
    with improved variant generation and content
    """
    
    def __init__(self, use_ai_generation=False, ai_model_path=None):
        """
        Initialize the enhanced slide generator
        
        Args:
            use_ai_generation: Whether to use AI for content generation
            ai_model_path: Path to the AI model (if applicable)
        """
        self.use_ai_generation = use_ai_generation
        self.ai_model_path = ai_model_path
        
        # Initialize AI model if needed
        if self.use_ai_generation and self.ai_model_path:
            try:
                # This would be the code to load a quantized model like DeepSeek v3
                # For now, we'll use template-based generation as a fallback
                pass
            except Exception as e:
                print(f"Error loading AI model: {e}")
                self.use_ai_generation = False
                
        # Templates for hypothesis generation (non-AI approach)
        self.hypothesis_templates = {
            "price": "We believe that adjusting pricing for {segment} will {expected_outcome} because {rationale}.",
            "shipping": "We believe that modifying shipping options for {segment} will {expected_outcome} because {rationale}.",
            "layout": "We believe that updating the interface design for {segment} will {expected_outcome} because {rationale}.",
            "messaging": "We believe that changing the messaging for {segment} will {expected_outcome} because {rationale}.",
            "generic": "We believe that the proposed changes will {expected_outcome} for {segment} because {rationale}."
        }
        
        self.outcome_templates = {
            "price": [
                "increase revenue without significantly impacting conversion rate",
                "optimize revenue while maintaining customer satisfaction",
                "improve profitability without affecting customer acquisition"
            ],
            "shipping": [
                "improve conversion rates",
                "reduce cart abandonment",
                "increase customer satisfaction"
            ],
            "layout": [
                "improve engagement and conversions",
                "reduce friction and increase conversion rate",
                "enhance user experience and increase time on page"
            ],
            "messaging": [
                "increase click-through rates",
                "improve conversion by resonating better with customers",
                "enhance brand perception and customer trust"
            ],
            "generic": [
                "positively impact performance",
                "improve key metrics",
                "drive better business outcomes"
            ]
        }
        
        self.rationale_templates = {
            "price": [
                "customers in this segment demonstrate high purchase intent",
                "we've observed price elasticity in this customer segment",
                "similar segments have shown tolerance to price adjustments"
            ],
            "shipping": [
                "these customers are sensitive to shipping costs and delivery times",
                "shipping options significantly influence purchase decisions for this segment",
                "delivery expectations for this segment align with our proposed changes"
            ],
            "layout": [
                "the new layout addresses common friction points identified in user feedback",
                "heatmap analysis indicates the current design causes user confusion",
                "similar layout changes have shown success with comparable segments"
            ],
            "messaging": [
                "our customer research indicates stronger alignment with this messaging approach",
                "this segment responds better to this communication style",
                "A/B tests in similar contexts have shown improved engagement with this type of messaging"
            ],
            "generic": [
                "previous customer behavior patterns suggest this change will be well-received",
                "market trends and competitive analysis support this approach",
                "our data indicates this is an area of opportunity"
            ]
        }

    def generate_hypothesis(self, test_type, segment, supporting_data_text=""):
        """
        Generate a hypothesis based on test type and segment
        
        Args:
            test_type: String describing the test
            segment: String describing the target segment
            supporting_data_text: Optional text from supporting data
            
        Returns:
            Generated hypothesis string
        """
        if self.use_ai_generation:
            # This would be the code to use AI for hypothesis generation
            # For now, we'll use template-based generation
            pass
            
        # Template-based generation (fallback)
        # Determine the test category
        category = "generic"
        test_type_lower = test_type.lower()
        
        if any(word in test_type_lower for word in ["price", "$", "cost"]):
            category = "price"
        elif any(word in test_type_lower for word in ["shipping", "delivery"]):
            category = "shipping"
        elif any(word in test_type_lower for word in ["layout", "design", "ui", "interface"]):
            category = "layout"
        elif any(word in test_type_lower for word in ["message", "copy", "text", "wording"]):
            category = "messaging"
        
        # Select random outcome and rationale for variety
        expected_outcome = random.choice(self.outcome_templates[category])
        rationale = random.choice(self.rationale_templates[category])
        
        # Format the template
        hypothesis = self.hypothesis_templates[category].format(
            segment=segment,
            expected_outcome=expected_outcome,
            rationale=rationale
        )
        
        return hypothesis
    
    def infer_goals_and_kpis(self, test_type):
        """
        Determine appropriate goals and KPIs based on test description
        
        Args:
            test_type: String describing the test
            
        Returns:
            Tuple of (goal, kpi, targeted_impact)
        """
        test_type_lower = test_type.lower()
        
        # Default values
        goal = "Improve Performance"
        kpi = "Conversion Rate (CVR)"
        targeted_impact = "3-5%"
        
        # Price-related tests
        if any(term in test_type_lower for term in ["price", "pricing", "cost", "$"]):
            goal = "Increase Revenue"
            kpi = "Revenue Per Visitor (RPV)"
            targeted_impact = "5-8%"
            
            # Check for specific price changes
            price_matches = re.findall(r'\$(\d+\.\d{2}|\d+)', test_type)
            if len(price_matches) >= 2:
                try:
                    # Calculate percentage difference between prices
                    price1 = float(price_matches[0])
                    price2 = float(price_matches[1])
                    if price1 > 0:
                        pct_diff = abs((price2 - price1) / price1) * 100
                        if pct_diff > 20:
                            targeted_impact = "8-12%"
                        elif pct_diff > 10:
                            targeted_impact = "5-8%"
                        else:
                            targeted_impact = "3-5%"
                except ValueError:
                    pass
        
        # Conversion-focused tests
        elif any(term in test_type_lower for term in ["conversion", "cvr", "checkout"]):
            goal = "Increase Conversion Rate"
            kpi = "Conversion Rate (CVR)"
            targeted_impact = "3-5%"
        
        # AOV-focused tests
        elif any(term in test_type_lower for term in ["aov", "order value", "cart"]):
            goal = "Increase Average Order Value"
            kpi = "Average Order Value (AOV)"
            targeted_impact = "7-10%"
        
        # Shipping/delivery tests
        elif any(term in test_type_lower for term in ["shipping", "delivery", "ship"]):
            goal = "Reduce Cart Abandonment"
            kpi = "Checkout Completion Rate"
            targeted_impact = "4-7%"
        
        # Layout/design tests
        elif any(term in test_type_lower for term in ["layout", "design", "ui", "interface"]):
            goal = "Improve User Experience"
            kpi = "Engagement Rate"
            targeted_impact = "5-10%"
        
        # Messaging/copy tests
        elif any(term in test_type_lower for term in ["message", "copy", "text", "wording", "note"]):
            goal = "Increase Engagement"
            kpi = "Click-Through Rate (CTR)"
            targeted_impact = "8-15%"
            
        return goal, kpi, targeted_impact
    
    def generate_tags(self, test_type, supporting_data_text=""):
        """
        Generate relevant tags based on test description and supporting data
        
        Args:
            test_type: String describing the test
            supporting_data_text: Optional text from supporting data
            
        Returns:
            List of tags
        """
        tags = []
        test_type_lower = test_type.lower()
        
        # Pattern matching for common test categories
        if any(word in test_type_lower for word in ["price", "$", "cost"]):
            tags.append("Price Sensitivity")
            
        if any(word in test_type_lower for word in ["shipping", "delivery", "ship"]):
            tags.append("Shipping Options")
            
        if any(word in test_type_lower for word in ["premium", "bundle", "upgrade"]):
            tags.append("Premium Offering")
            
        if any(word in test_type_lower for word in ["layout", "design", "ui", "ux"]):
            tags.append("UI/UX Design")
            
        if any(word in test_type_lower for word in ["message", "copy", "text", "wording", "note"]):
            tags.append("Messaging")
            
        if any(word in test_type_lower for word in ["checkout", "cart", "payment"]):
            tags.append("Checkout Process")
            
        if "free" in test_type_lower:
            tags.append("Free Offer")
            
        if "discount" in test_type_lower:
            tags.append("Discount")
            
        if "mobile" in test_type_lower:
            tags.append("Mobile Experience")
            
        if "desktop" in test_type_lower:
            tags.append("Desktop Experience")
        
        if "time" in test_type_lower or "pm" in test_type_lower or "am" in test_type_lower:
            tags.append("Time-Based")
            
        if "dynamic" in test_type_lower:
            tags.append("Dynamic Content")
            
        # Use supporting data text to identify additional tags if available
        combined_text = f"{test_type} {supporting_data_text}".lower()
        
        if any(term in combined_text for term in ["first time", "new user", "new customer"]):
            tags.append("New Customer")
            
        if any(term in combined_text for term in ["mobile", "smartphone", "android", "ios"]):
            tags.append("Mobile Experience")
            
        if any(term in combined_text for term in ["returning", "repeat", "loyal"]):
            tags.append("Returning Customer")
            
        if any(term in combined_text for term in ["abandoned", "cart abandonment"]):
            tags.append("Cart Abandonment")
            
        if any(term in combined_text for term in ["fst", "free shipping threshold"]):
            tags.append("Free Shipping Threshold")
            
        # Limit to top 3 most relevant tags
        return list(set(tags))[:3]
    
    def detect_price_text(self, image, old_price=None):
        """
        Enhanced method to detect price text in an image using OCR and pattern matching
        
        Args:
            image: PIL Image
            old_price: Optional price text to find (e.g., "$3.00")
            
        Returns:
            List of (x, y, width, height) locations
        """
        # Convert PIL to OpenCV format
        img_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        
        # Convert to grayscale
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
        
        # Apply thresholding
        _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY_INV)
        
        # Perform OCR with detailed information including bounding boxes
        ocr_data = pytesseract.image_to_data(thresh, output_type=pytesseract.Output.DICT)
        
        # Find price locations
        price_locations = []
        
        # Price pattern in various formats
        price_pattern = r'\$\d+\.\d{2}|\$\d+'
        
        # Process OCR data
        for i in range(len(ocr_data['text'])):
            text = ocr_data['text'][i]
            
            # Check if text contains a price
            if re.search(price_pattern, text):
                # If specific price is provided, check for a match
                if old_price and old_price not in text:
                    continue
                    
                x = ocr_data['left'][i]
                y = ocr_data['top'][i]
                w = ocr_data['width'][i]
                h = ocr_data['height'][i]
                
                # Store location with the detected price
                detected_price = re.search(price_pattern, text).group(0)
                price_locations.append({
                    "x": x,
                    "y": y, 
                    "width": w,
                    "height": h,
                    "detected_price": detected_price
                })
        
        # If we couldn't find it with OCR, use a more aggressive approach
        if not price_locations and old_price:
            # Search specifically for the old_price using template matching
            # This is a placeholder for more sophisticated methods
            print(f"OCR could not detect price '{old_price}'. Using fallback method.")
            
            # Extract price digits for digit-based detection
            price_digits = re.sub(r'[^\d.]', '', old_price)
            if price_digits:
                for i in range(len(ocr_data['text'])):
                    text = ocr_data['text'][i]
                    if price_digits in text:
                        x = ocr_data['left'][i]
                        y = ocr_data['top'][i]
                        w = ocr_data['width'][i]
                        h = ocr_data['height'][i]
                        
                        price_locations.append({
                            "x": x,
                            "y": y, 
                            "width": w,
                            "height": h,
                            "detected_price": old_price
                        })
        
        # If still no locations found, use a fixed position
        if not price_locations:
            # Check for price in standard shipping row - hardcoded positioned based on shipping template
            h, w = img_cv.shape[:2]
            
            # Typical position for standard shipping price in checkout forms
            right_side_x = int(w * 0.85)  # 85% of the way to the right
            first_row_y = int(h * 0.25)   # 25% of the way down
            second_row_y = int(h * 0.6)   # 60% of the way down
            
            # Add both potential positions - first is for Standard shipping, second is for Rush
            price_locations.append({
                "x": right_side_x - 60,
                "y": first_row_y - 10,
                "width": 60,
                "height": 30,
                "detected_price": old_price or "$7.95"
            })
            
            price_locations.append({
                "x": right_side_x - 60,
                "y": second_row_y - 10,
                "width": 60,
                "height": 30,
                "detected_price": old_price or "$24.95" 
            })
            
        return price_locations
    
    def create_price_variant_mockup(self, control_image, test_type):
        """
        Create a variant mockup with price changes - with enhanced precision
        
        Args:
            control_image: PIL Image of control
            test_type: Test description to extract price information
            
        Returns:
            PIL Image with price modifications
        """
        # Extract prices from test type
        old_price = None
        new_price = None
        
        prices = re.findall(r'\$\d+\.\d{2}|\$\d+', test_type)
        if len(prices) >= 2:
            old_price, new_price = prices[0], prices[1]
        else:
            # If no specific prices found, use defaults for shipping tests
            old_price = "$7.95"
            # Use a price reduction for standard shipping
            new_price = "$5.95"
        
        # Create a copy of the control image
        variant_image = control_image.copy()
        variant_array = np.array(variant_image)
        
        # Convert to cv2 format for processing
        img_cv = cv2.cvtColor(variant_array, cv2.COLOR_RGB2BGR)
        
        # Attempt to detect price locations
        price_locations = self.detect_price_text(control_image, old_price)
        
        # If we found price locations, modify them with the new price
        for location in price_locations:
            x = location["x"]
            y = location["y"]
            w = location["width"]
            h = location["height"]
            
            # Create a clean rectangle to cover old price
            cv2.rectangle(img_cv, (x-5, y-5), (x+w+5, y+h+5), (255, 255, 255), -1)
            
            # Add new price text
            font = cv2.FONT_HERSHEY_DUPLEX
            font_scale = 0.6
            font_thickness = 1
            
            # Draw the new price - slightly different position for better alignment
            cv2.putText(img_cv, new_price, (x, y+h-3), 
                      font, font_scale, (0, 0, 0), font_thickness, cv2.LINE_AA)
        
        # Add a "VARIANT" indicator in the top-right corner
        h, w = img_cv.shape[:2]
        
        # Create a banner for "VARIANT" text
        banner_width = 100
        banner_height = 30
        banner_x = w - banner_width - 10  # 10px from right edge
        banner_y = 10  # 10px from top
        
        # Background for variant indicator
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (255, 255, 255), -1)  # White background
        
        # Add border
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (255, 0, 0), 2)  # Red border
        
        # Add "VARIANT" text
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.5
        text_size = cv2.getTextSize("VARIANT", font, font_scale, 2)[0]
        text_x = banner_x + (banner_width - text_size[0]) // 2
        text_y = banner_y + banner_height - 10
        
        cv2.putText(img_cv, "VARIANT", (text_x, text_y), 
                  font, font_scale, (255, 0, 0), 2, cv2.LINE_AA)
        
        # Convert back to PIL
        return Image.fromarray(cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB))
    
    def create_shipping_variant_mockup(self, control_image, test_description):
        """
        Create a variant mockup for shipping tests
        
        Args:
            control_image: PIL Image of control
            test_description: Description of the shipping test
            
        Returns:
            PIL Image with shipping modifications
        """
        # Create a copy of the control image
        variant_image = control_image.copy()
        variant_array = np.array(variant_image)
        
        # Convert to cv2 format
        img_cv = cv2.cvtColor(variant_array, cv2.COLOR_RGB2BGR)
        
        # Determine the type of shipping test
        test_description_lower = test_description.lower()
        
        if "before" in test_description_lower and "after" in test_description_lower and "pm" in test_description_lower:
            # Dynamic time-based messaging test
            message = "Order ships tomorrow if placed after 4 PM"
            color = (0, 0, 255)  # Red text
        elif "expedited" in test_description_lower or "express" in test_description_lower:
            # Expedited shipping test
            message = "Expedited Shipping Available"
            color = (0, 128, 0)  # Green text
        elif "free" in test_description_lower:
            # Free shipping test
            message = "Free Shipping Eligible"
            color = (0, 128, 0)  # Green text
        else:
            # Generic shipping change
            message = "Updated Shipping Option"
            color = (0, 0, 0)  # Black text
        
        # Add banner with shipping message
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.7
        text_size = cv2.getTextSize(message, font, font_scale, 2)[0]
        
        # Position banner at the top of the image
        h, w = img_cv.shape[:2]
        banner_x = 10
        banner_y = 30
        
        # Add a background highlight
        cv2.rectangle(img_cv, 
                    (banner_x - 5, banner_y - text_size[1] - 5),
                    (banner_x + text_size[0] + 15, banner_y + 5),
                    (255, 255, 255), -1)
        
        # Add a colored border
        cv2.rectangle(img_cv, 
                    (banner_x - 5, banner_y - text_size[1] - 5),
                    (banner_x + text_size[0] + 15, banner_y + 5),
                    color, 2)
        
        # Add the banner text
        cv2.putText(img_cv, message, (banner_x, banner_y), 
                  font, font_scale, color, 2, cv2.LINE_AA)
        
        # Add a "VARIANT" indicator in the top-right corner
        banner_width = 100
        banner_height = 30
        banner_x = w - banner_width - 10  # 10px from right edge
        banner_y = 10  # 10px from top
        
        # Background for variant indicator
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (255, 255, 255), -1)  # White background
        
        # Add border
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (0, 0, 255), 2)  # Red border
        
        # Add "VARIANT" text
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.5
        text_size = cv2.getTextSize("VARIANT", font, font_scale, 2)[0]
        text_x = banner_x + (banner_width - text_size[0]) // 2
        text_y = banner_y + banner_height - 10
        
        cv2.putText(img_cv, "VARIANT", (text_x, text_y), 
                  font, font_scale, (255, 0, 0), 2, cv2.LINE_AA)
        
        # Convert back to PIL
        return Image.fromarray(cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB))
    
    def create_messaging_variant_mockup(self, control_image, test_description):
        """
        Create a variant mockup for messaging/copy tests
        
        Args:
            control_image: PIL Image of control
            test_description: Description of the messaging test
            
        Returns:
            PIL Image with messaging modifications
        """
        # Create a copy of the control image
        variant_image = control_image.copy()
        variant_array = np.array(variant_image)
        
        # Convert to cv2 format
        img_cv = cv2.cvtColor(variant_array, cv2.COLOR_RGB2BGR)
        
        # Determine the type of messaging test
        test_description_lower = test_description.lower()
        
        # Select appropriate messaging based on test description
        if "urgency" in test_description_lower:
            message = "Limited Time Offer!"
            color = (0, 0, 255)  # Red for urgency
        elif "social proof" in test_description_lower:
            message = "Bestseller - Customer Favorite"
            color = (0, 128, 0)  # Green for positive reinforcement
        elif "note" in test_description_lower and ("before" in test_description_lower or "after" in test_description_lower):
            if "before" in test_description_lower and "after" in test_description_lower:
                # For before/after tests, create a split message
                message = "Order TODAY for Same Day Shipping!"
                color = (0, 0, 255)  # Red for urgency
            else:
                message = "Updated Shipping Message"
                color = (0, 0, 255)
        else:
            message = "New Messaging Variant"
            color = (0, 0, 0)  # Black for general
        
        # Add banner with messaging
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.7
        text_size = cv2.getTextSize(message, font, font_scale, 2)[0]
        
        # Position banner at the top of the image
        h, w = img_cv.shape[:2]
        banner_x = 10
        banner_y = 30
        
        # Add a background highlight
        cv2.rectangle(img_cv, 
                    (banner_x - 5, banner_y - text_size[1] - 5),
                    (banner_x + text_size[0] + 15, banner_y + 5),
                    (255, 255, 255), -1)
        
        # Add a colored border
        cv2.rectangle(img_cv, 
                    (banner_x - 5, banner_y - text_size[1] - 5),
                    (banner_x + text_size[0] + 15, banner_y + 5),
                    color, 2)
        
        # Add the banner text
        cv2.putText(img_cv, message, (banner_x, banner_y), 
                  font, font_scale, color, 2, cv2.LINE_AA)
        
        # For before/after tests, add a specific indicator for time-based changes
        if "before" in test_description_lower and "after" in test_description_lower:
            time_message = "Message changes after 4PM"
            
            # Add time indicator at bottom of image
            time_banner_x = 10
            time_banner_y = h - 20
            time_text_size = cv2.getTextSize(time_message, font, 0.6, 1)[0]
            
            # Add background for time message
            cv2.rectangle(img_cv, 
                        (time_banner_x - 5, time_banner_y - time_text_size[1] - 5),
                        (time_banner_x + time_text_size[0] + 15, time_banner_y + 5),
                        (255, 255, 255), -1)
            
            # Add time message
            cv2.putText(img_cv, time_message, (time_banner_x, time_banner_y), 
                      font, 0.6, (0, 0, 0), 1, cv2.LINE_AA)
        
        # Add a "VARIANT" indicator in the top-right corner
        banner_width = 100
        banner_height = 30
        banner_x = w - banner_width - 10  # 10px from right edge
        banner_y = 10  # 10px from top
        
        # Background for variant indicator
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (255, 255, 255), -1)  # White background
        
        # Add border
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (0, 0, 255), 2)  # Red border
        
        # Add "VARIANT" text
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.5
        text_size = cv2.getTextSize("VARIANT", font, font_scale, 2)[0]
        text_x = banner_x + (banner_width - text_size[0]) // 2
        text_y = banner_y + banner_height - 10
        
        cv2.putText(img_cv, "VARIANT", (text_x, text_y), 
                  font, font_scale, (255, 0, 0), 2, cv2.LINE_AA)
        
        # Convert back to PIL
        return Image.fromarray(cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB))
    
    def create_layout_variant_mockup(self, control_image, test_description):
        """
        Create a variant mockup for layout/design tests
        
        Args:
            control_image: PIL Image of control
            test_description: Description of the layout test
            
        Returns:
            PIL Image with layout modifications
        """
        # Create a copy of the control image
        variant_image = control_image.copy()
        variant_array = np.array(variant_image)
        
        # Convert to cv2 format
        img_cv = cv2.cvtColor(variant_array, cv2.COLOR_RGB2BGR)
        
        # Get dimensions
        h, w = img_cv.shape[:2]
        
        # Add a layout indicator
        layout_text = "Modified Layout"
        font = cv2.FONT_HERSHEY_SIMPLEX
        
        # Add highlight region to indicate layout change
        # Draw a rectangle around the central portion of the image
        margin_x = int(w * 0.1)
        margin_y = int(h * 0.15)
        
        # Draw rectangle to highlight the changed region
        cv2.rectangle(img_cv,
                    (margin_x, margin_y),
                    (w - margin_x, h - margin_y),
                    (0, 165, 255), 3)  # Orange border
        
        # Add layout indicator text
        cv2.putText(img_cv, layout_text, (margin_x + 10, margin_y - 10),
                  font, 0.7, (0, 165, 255), 2, cv2.LINE_AA)
        
        # Add a "VARIANT" indicator in the top-right corner
        banner_width = 100
        banner_height = 30
        banner_x = w - banner_width - 10  # 10px from right edge
        banner_y = 10  # 10px from top
        
        # Background for variant indicator
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (255, 255, 255), -1)  # White background
        
        # Add border
        cv2.rectangle(img_cv, 
                    (banner_x, banner_y),
                    (banner_x + banner_width, banner_y + banner_height),
                    (0, 0, 255), 2)  # Red border
        
        # Add "VARIANT" text
        font = cv2.FONT_HERSHEY_SIMPLEX
        font_scale = 0.5
        text_size = cv2.getTextSize("VARIANT", font, font_scale, 2)[0]
        text_x = banner_x + (banner_width - text_size[0]) // 2
        text_y = banner_y + banner_height - 10
        
        cv2.putText(img_cv, "VARIANT", (text_x, text_y), 
                  font, font_scale, (255, 0, 0), 2, cv2.LINE_AA)
        
        # Convert back to PIL
        return Image.fromarray(cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB))
        
    def create_variant_mockup(self, control_image, test_type):
        """
        Create an appropriate variant mockup based on test type
        
        Args:
            control_image: PIL Image of control
            test_type: Test description
            
        Returns:
            PIL Image of variant
        """
        test_type_lower = test_type.lower()
        
        # Price/cost test
        if any(word in test_type_lower for word in ["price", "$", "cost"]):
            return self.create_price_variant_mockup(control_image, test_type)
        
        # Shipping/delivery test
        elif any(word in test_type_lower for word in ["shipping", "delivery", "ship"]):
            return self.create_shipping_variant_mockup(control_image, test_type)
        
        # Messaging/copy test
        elif any(word in test_type_lower for word in ["message", "copy", "text", "wording", "note"]):
            return self.create_messaging_variant_mockup(control_image, test_type)
        
        # Layout/design test
        elif any(word in test_type_lower for word in ["layout", "design", "ui", "ux"]):
            return self.create_layout_variant_mockup(control_image, test_type)
        
        # Default variant - generic modification with price change
        # For any test not clearly categorized, default to price sensitivity
        return self.create_price_variant_mockup(control_image, test_type)
    
    def create_pptx_slide(self, slide_data):
        """
        Generate a PowerPoint slide with the test information using the modern purple design
        
        Args:
            slide_data: Dictionary containing all slide information
            
        Returns:
            BytesIO buffer containing the PPTX
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Add a blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Set dark purple background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(32, 21, 54)  # Dark purple background (#201536)
        
        # Add small logo in top-left corner if available
        try:
            logo_path = os.path.join(os.path.dirname(__file__), "logo.png")
            if os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(0.6))
        except:
            pass  # Skip logo if not available
        
        # Add title at top
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
        title_frame = title.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = slide_data['title']
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # SECTION 1: Hypothesis (purple box)
        hypothesis_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.0), Inches(1.2))
        hypothesis_box.fill.solid()
        hypothesis_box.fill.fore_color.rgb = RGBColor(48, 34, 76)  # Lighter purple
        
        hypothesis_frame = hypothesis_box.text_frame
        hypothesis_frame.word_wrap = True
        
        # Add pencil icon for hypothesis
        pencil_icon = "✏️"
        hypothesis_para = hypothesis_frame.add_paragraph()
        hypothesis_para.text = f"{pencil_icon} Hypothesis"
        hypothesis_para.font.size = Pt(14)
        hypothesis_para.font.bold = True
        hypothesis_para.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # Add hypothesis text
        hypothesis_text = hypothesis_frame.add_paragraph()
        hypothesis_text.text = slide_data['hypothesis']
        hypothesis_text.font.size = Pt(12)
        hypothesis_text.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # SECTION 2: Create a 2x2 grid of info boxes under the hypothesis
        grid_top = Inches(3.0)
        grid_left = Inches(0.5)
        box_width = Inches(3.5)
        box_height = Inches(1.0)
        
        # Grid box styling
        def create_info_box(left, top, width, height, title, content, icon=""):
            box = slide.shapes.add_textbox(left, top, width, height)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(48, 34, 76)  # Lighter purple
            
            text_frame = box.text_frame
            text_frame.word_wrap = True
            
            # Title with icon
            title_para = text_frame.add_paragraph()
            title_para.text = f"{icon} {title}" if icon else title
            title_para.font.size = Pt(14)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
            
            # Content
            content_para = text_frame.add_paragraph()
            content_para.text = content
            content_para.font.size = Pt(12)
            content_para.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # Create the grid boxes with their respective icons and content
        # Box 1: Segment
        create_info_box(
            grid_left, 
            grid_top, 
            box_width, 
            box_height, 
            "Segment", 
            slide_data['segment'],
            "👥"
        )
        
        # Box 2: Timeline
        create_info_box(
            grid_left + box_width + Inches(0.25), 
            grid_top, 
            box_width, 
            box_height, 
            "Timeline", 
            f"{slide_data.get('timeline', '2-3 weeks')}\nStat Sig: {slide_data.get('stat_sig', '85%')}",
            "📅"
        )
        
        # Box 3: Elements (tags)
        create_info_box(
            grid_left, 
            grid_top + box_height + Inches(0.25), 
            box_width, 
            box_height, 
            "Elements", 
            ", ".join(slide_data['tags'])
        )
        
        # Box 4: Success Criteria
        create_info_box(
            grid_left + box_width + Inches(0.25), 
            grid_top + box_height + Inches(0.25), 
            box_width, 
            box_height, 
            "Success Criteria", 
            f"{slide_data.get('success_criteria', 'Increase in Revenue without significant drop in Conversion Rate')}",
            "✓"
        )
        
        # Box 5: Goal (NEW)
        create_info_box(
            grid_left,
            grid_top + box_height*2 + Inches(0.5),
            box_width,
            box_height,
            "Goal",
            slide_data['goal'],
            "🎯"
        )
        
        # Box 6: KPI and Targeted Impact (NEW)
        create_info_box(
            grid_left + box_width + Inches(0.25),
            grid_top + box_height*2 + Inches(0.5),
            box_width,
            box_height,
            "KPI & Targeted Impact",
            f"{slide_data['kpi']} ({slide_data['impact']} improvement)"
        )
        
        # SECTION 3: Variant box (positioned in the top right)
        variant_box = slide.shapes.add_textbox(
            Inches(8.0), 
            Inches(1.5), 
            Inches(4.33), 
            Inches(0.4)  # Shorter box for just the title
        )
        variant_box.fill.solid()
        variant_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
        variant_frame = variant_box.text_frame
        variant_title = variant_frame.add_paragraph()
        
        # Special handling for dynamic messaging tests
        if "before" in slide_data['title'].lower() and "after" in slide_data['title'].lower():
            variant_title.text = "Variant B - Before/After 4PM"
        else:
            variant_title.text = "Variant"
            
        variant_title.font.size = Pt(16)
        variant_title.font.bold = True
        
        # Save variant image to temp file and add to slide
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as variant_temp:
            variant_img = slide_data['variant_image']
            variant_img.save(variant_temp, format='PNG')
            variant_temp_path = variant_temp.name
        
        # Add variant image with proper positioning
        slide.shapes.add_picture(
            variant_temp_path, 
            Inches(8.0), 
            Inches(2.0),  # Positioned directly below the title box
            width=Inches(4.33)
        )
        
        # SECTION 4: Control box (positioned below variant)
        control_box = slide.shapes.add_textbox(
            Inches(8.0), 
            Inches(4.75), 
            Inches(4.33), 
            Inches(0.4)  # Shorter box for just the title
        )
        control_box.fill.solid()
        control_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
        control_frame = control_box.text_frame
        control_title = control_frame.add_paragraph()
        control_title.text = "Control"
        control_title.font.size = Pt(16)
        control_title.font.bold = True
        
        # Save control image to temp file and add to slide
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as control_temp:
            control_img = slide_data['control_image']
            control_img.save(control_temp, format='PNG')
            control_temp_path = control_temp.name
        
        # Add control image with proper positioning
        slide.shapes.add_picture(
            control_temp_path, 
            Inches(8.0), 
            Inches(5.2),  # Positioned directly below the title box
            width=Inches(4.33)
        )
        
        # SECTION 5: Supporting Data (at the bottom left)
        if slide_data.get('supporting_data_image'):
            # Save supporting data image to temp file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as supporting_temp:
                supporting_img = slide_data['supporting_data_image']
                supporting_img.save(supporting_temp, format='PNG')
                supporting_temp_path = supporting_temp.name
            
            # Add supporting data image
            slide.shapes.add_picture(
                supporting_temp_path,
                Inches(0.5),
                Inches(5.5),
                width=Inches(6.0)
            )
        
        # SECTION 6: Required Checkouts Indicator (NEW)
        checkout_box = slide.shapes.add_textbox(
            Inches(3.0),
            Inches(5.5),
            Inches(2.0),
            Inches(1.0)
        )
        checkout_frame = checkout_box.text_frame
        checkout_para = checkout_frame.add_paragraph()
        checkout_para.text = "🛍️ 20,000"
        checkout_para.font.size = Pt(24)
        checkout_para.font.bold = True
        checkout_para.alignment = PP_ALIGN.CENTER
        checkout_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # Add description for checkout number
        checkout_desc = checkout_frame.add_paragraph()
        checkout_desc.text = "# of checkouts required"
        checkout_desc.font.size = Pt(12)
        checkout_desc.alignment = PP_ALIGN.CENTER
        checkout_desc.font.color.rgb = RGBColor(255, 255, 255)  # White text
            
        # Save presentation to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        
        # Clean up temporary files
        os.unlink(control_temp_path)
        os.unlink(variant_temp_path)
        if slide_data.get('supporting_data_image'):
            os.unlink(supporting_temp_path)
        
        buffer.seek(0)
        return buffer
        
    def determine_success_criteria(self, test_type, kpi):
        """
        Determine appropriate success criteria based on test type and KPI
        
        Args:
            test_type: String describing the test
            kpi: Key Performance Indicator for the test
            
        Returns:
            String describing success criteria
        """
        test_type_lower = test_type.lower()
        
        # Default success criteria
        criteria = f"Significant improvement in {kpi}"
        
        # Price tests
        if any(term in test_type_lower for term in ["price", "pricing", "cost", "$"]):
            criteria = "Increase in Revenue without significant drop in Conversion Rate"
            
        # Conversion tests
        elif any(term in test_type_lower for term in ["conversion", "cvr", "checkout"]):
            criteria = "Uplift of 1%-2% in conversion rate"
            
        # AOV-focused tests
        elif any(term in test_type_lower for term in ["aov", "order value"]):
            criteria = "Increase in Average Order Value without drop in Conversion Rate"
            
        # Shipping tests
        elif any(term in test_type_lower for term in ["shipping", "delivery", "ship"]):
            criteria = "Decrease in Cart Abandonment Rate"
            
        # Layout/design tests
        elif any(term in test_type_lower for term in ["layout", "design", "ui"]):
            criteria = "Improvement in User Engagement Metrics"
            
        # Messaging/copy tests
        elif any(term in test_type_lower for term in ["message", "copy", "text", "note"]):
            criteria = "Increase in Click-Through Rate"
            
        return criteria
    
    def parse_test_type(self, test_type):
        """
        Extract a clean title from the test type description
        
        Args:
            test_type: Full test type description
            
        Returns:
            Clean title string
        """
        # Look for a structured format like "Test Type — Control: X | Variant: Y"
        parts = test_type.split('—')
        
        if len(parts) > 1:
            return parts[0].strip()
        
        # Look for the pattern "Test: Description"
        colon_split = test_type.split(':')
        if len(colon_split) > 1:
            return colon_split[0].strip()
        
        # If no structured format, return as is or truncate if too long
        if len(test_type) > 40:
            return test_type[:37] + "..."
        
        return test_type
    
    def generate_slide(self, test_type, segment, control_image, supporting_data_image=None, 
                      supporting_data_text="", custom_hypothesis=None):
        """
        Generate a complete A/B test slide
        
        Args:
            test_type: String describing the test
            segment: String describing the target segment
            control_image: PIL Image of the control layout
            supporting_data_image: Optional PIL Image with supporting data
            supporting_data_text: Optional text extracted from supporting data
            custom_hypothesis: Optional custom hypothesis
            
        Returns:
            BytesIO buffer containing the PPTX slide
        """
        # Parse test type
        title = f"A/B Test: {self.parse_test_type(test_type)}"
        
        # Generate hypothesis
        if custom_hypothesis:
            hypothesis = custom_hypothesis
        else:
            hypothesis = self.generate_hypothesis(test_type, segment, supporting_data_text)
        
        # Determine goal, KPI, impact
        goal, kpi, impact = self.infer_goals_and_kpis(test_type)
        
        # Generate tags
        tags = self.generate_tags(test_type, supporting_data_text)
        
        # Generate variant mockup - this now uses enhanced processing
        variant_image = self.create_variant_mockup(control_image, test_type)
        
        # Determine success criteria
        success_criteria = self.determine_success_criteria(test_type, kpi)
        
        # Create slide data dictionary
        slide_data = {
            "title": title,
            "hypothesis": hypothesis,
            "segment": segment,
            "goal": goal,
            "kpi": kpi,
            "impact": impact,
            "tags": tags,
            "control_image": control_image,
            "variant_image": variant_image,
            "supporting_data_image": supporting_data_image,
            "success_criteria": success_criteria,
            "timeline": "2-3 weeks",
            "stat_sig": "85%"
        }
        
        # Generate PowerPoint slide
        return self.create_pptx_slide(slide_data)