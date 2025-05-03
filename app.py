# -*- coding: utf-8 -*-
import streamlit as st
import cv2
import numpy as np
import pytesseract # Keep this import for the python wrapper
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont, UnidentifiedImageError
import io
import os
import base64
import random
import tempfile
import re
import logging
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR # Correct import
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime
from html2image import Html2Image # Import html2image

# Initialize session state variables
if "slide_generated" not in st.session_state:
    st.session_state.slide_generated = False
if "output_buffer" not in st.session_state:
    st.session_state.output_buffer = None
if "slide_data" not in st.session_state:
    st.session_state.slide_data = {}
if "error_message" not in st.session_state:
    st.session_state.error_message = None
if "control_image_filename" not in st.session_state:
    st.session_state.control_image_filename = None
if "control_image_data" not in st.session_state:
    st.session_state.control_image_data = None

# --- Basic Setup ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(name)s - %(message)s')
logger = logging.getLogger("ABTestApp")

# PDQ Color Palette
PDQ_COLORS = {
    "revolver": "#231333", "electric_violet": "#894DFF", "electric_violet_2": "#6B2BEF",
    "magnolia": "#F6F3FF", "melrose": "#9B8CFF", "valentino": "#2C124B",
    "moon_raker": "#E4E1FA", "lightning_yellow": "#FBC018", "persian_green": "#00AA8C",
    "carnation": "#F04557", "white": "#FFFFFF", "black": "#000000",
    "grey_text": "#666666", "html_border": "#DDDDDD", "html_selected_border": "#E86C60",
    "html_selected_bg": "#FEF6F5", "html_radio_border": "#CCCCCC",
}

# --- Helper Functions ---
def hex_to_rgb(hex_color):
    """Converts hex color string to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    length = len(hex_color)
    if length == 6: return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    if length == 3: return tuple(int(hex_color[i]*2, 16) for i in (0, 1, 2))
    logger.warning(f"Invalid hex: {hex_color}. Using black."); return (0, 0, 0)

def hex_to_rgbcolor(hex_color):
    """Converts hex color string to python-pptx RGBColor object."""
    r, g, b = hex_to_rgb(hex_color); return RGBColor(r, g, b)

def get_font(name, size, bold=False, italic=False):
    """Attempts to load Segoe UI, falls back to Arial/Calibri/Liberation, then default."""
    common_fonts = [name, 'Segoe UI', 'Arial', 'Calibri', 'Liberation Sans'] # Added Liberation Sans
    font_path = None
    for font_name in common_fonts:
        try:
            style_suffix = ""
            if bold and italic: style_suffix = " Bold Italic"
            elif bold: style_suffix = " Bold"
            elif italic: style_suffix = " Italic"
            for ext in ['.ttf', '.otf']:
                try:
                    # Try exact name + style
                    return ImageFont.truetype(f"{font_name}{style_suffix}{ext}", size)
                except IOError:
                    try:
                        # Try just name (for fonts where style is part of the base name)
                        return ImageFont.truetype(f"{font_name}{ext}", size)
                    except IOError:
                        continue # Try next extension or font name
        except IOError:
             continue # Try next font in list

    logger.warning(f"Could not find specified fonts: {common_fonts}. Using default PIL font.")
    try:
        return ImageFont.load_default(size=size) # Try specifying size for default
    except AttributeError: # Older PIL might not support size argument for load_default
        return ImageFont.load_default()

# --- Streamlit Page Setup & CSS ---
# DO NOT CHANGE - UI and Theme are fixed
st.set_page_config(page_title="PDQ A/B Test Slide Generator", page_icon="🧪", layout="wide")
# Define custom CSS
st.markdown(f"""
    <style>
    /* ... Existing CSS - DO NOT CHANGE ... */
    .main {{ background-color: {PDQ_COLORS["magnolia"]}; }}
    .stApp {{ max-width: 1400px; margin: 0 auto; }}
    .success-box {{ padding: 1rem; border-radius: 0.5rem; background-color: rgba({hex_to_rgb(PDQ_COLORS["persian_green"])[0]}, {hex_to_rgb(PDQ_COLORS["persian_green"])[1]}, {hex_to_rgb(PDQ_COLORS["persian_green"])[2]}, 0.2); color: {PDQ_COLORS["persian_green"]}; margin-bottom: 1rem; border: 1px solid {PDQ_COLORS["persian_green"]}; }}
    .stButton>button {{ background-color: {PDQ_COLORS["electric_violet"]}; color: white; font-weight: bold; border: none; padding: 0.6rem 1.2rem; border-radius: 0.3rem; }}
    .stButton>button:hover {{ background-color: {PDQ_COLORS["electric_violet_2"]}; }}
    .stButton>button:disabled {{ background-color: #cccccc; color: #666666; cursor: not-allowed; }}
    h1, h2, h3 {{ color: {PDQ_COLORS["valentino"]}; }}
    .preview-box {{ border: 1px solid {PDQ_COLORS["melrose"]}; border-radius: 0.5rem; padding: 1.5rem; background-color: {PDQ_COLORS["white"]}; box-shadow: 0 2px 4px rgba(0,0,0,0.1); }}
    .download-button {{ background-color: {PDQ_COLORS["persian_green"]}; color: white; padding: 0.6rem 1.2rem; border-radius: 0.3rem; text-decoration: none; font-weight: bold; display: inline-block; margin-top: 1rem; border: none; }}
    .download-button:hover {{ background-color: #008a70; }}
    .stImage > img {{ border: 1px solid {PDQ_COLORS["moon_raker"]}; border-radius: 0.25rem; }}
    .stSidebar > div:first-child {{ background-color: {PDQ_COLORS["revolver"]}; }}
    .stSidebar .stMarkdown p, .stSidebar .stFileUploader label, .stSidebar .stTextInput label, .stSidebar .stTextArea label, .stSidebar .stCheckbox label, .stSidebar h1, .stSidebar h2, .stSidebar h3, .stSidebar h4 {{ color: {PDQ_COLORS["magnolia"]}; }}
    .stSidebar .stButton>button {{ background-color: {PDQ_COLORS["melrose"]}; color: {PDQ_COLORS["revolver"]}; }}
    .stSidebar .stButton>button:hover {{ background-color: {PDQ_COLORS["electric_violet"]}; color: white; }}
    .stTabs [data-baseweb="tab-list"] {{ gap: 2px; }}
    .stTabs [data-baseweb="tab"] {{ background-color: {PDQ_COLORS["moon_raker"]}; border-radius: 4px 4px 0 0; color: {PDQ_COLORS["valentino"]}; padding: 0.5rem 1rem; font-weight: 500; }}
    .stTabs [aria-selected="true"] {{ background-color: {PDQ_COLORS["electric_violet"]}; color: white; }}
    .stTextInput>div>div>input, .stTextArea>div>textarea {{ border: 1px solid {PDQ_COLORS["melrose"]}; border-radius: 0.25rem; }}
    .stTextInput>div>div>input:focus, .stTextArea>div>textarea:focus {{ border-color: {PDQ_COLORS["electric_violet"]}; box-shadow: 0 0 0 2px rgba({hex_to_rgb(PDQ_COLORS["electric_violet"])[0]}, {hex_to_rgb(PDQ_COLORS["electric_violet"])[1]}, {hex_to_rgb(PDQ_COLORS["electric_violet"])[2]}, 0.3); }}
    label {{ color: {PDQ_COLORS["valentino"]}; font-weight: 500; }}
    .error-box {{ padding: 1rem; border-radius: 0.5rem; background-color: rgba({hex_to_rgb(PDQ_COLORS["carnation"])[0]}, {hex_to_rgb(PDQ_COLORS["carnation"])[1]}, {hex_to_rgb(PDQ_COLORS["carnation"])[2]}, 0.2); color: {PDQ_COLORS["carnation"]}; margin-bottom: 1rem; border: 1px solid {PDQ_COLORS["carnation"]}; }}
    .highlight-box {{ padding: 1rem; border-radius: 0.5rem; background-color: rgba({hex_to_rgb(PDQ_COLORS["lightning_yellow"])[0]}, {hex_to_rgb(PDQ_COLORS["lightning_yellow"])[1]}, {hex_to_rgb(PDQ_COLORS["lightning_yellow"])[2]}, 0.2); color: #805b00; margin-bottom: 1rem; border: 1px solid {PDQ_COLORS["lightning_yellow"]}; }}
    footer {{ visibility: hidden; }}
    .custom-footer {{ margin-top: 2rem; padding-top: 1rem; border-top: 1px solid {PDQ_COLORS["moon_raker"]}; display: flex; justify-content: space-between; align-items: center; font-size: 0.85rem; }}
    .footer-left {{ color: {PDQ_COLORS["valentino"]}; }}
    .footer-right {{ color: {PDQ_COLORS["electric_violet"]}; font-weight: bold; }}
    </style>
""", unsafe_allow_html=True)


# --- Core Image/Text Processing Functions ---

# MODIFIED extract_text_from_image function to be more robust
def extract_text_from_image(image_input):
    """Extract text from an image file or PIL object using OCR"""
    try:
        image = None # Initialize image variable
        input_type = "Unknown"
        input_details = ""

        if isinstance(image_input, Image.Image):
            image = image_input
            input_type = "PIL Image"
            input_details = f"size={image.size}, mode={image.mode}"
            logger.info("Processing pre-loaded PIL image for OCR.")

        elif hasattr(image_input, 'getvalue') and hasattr(image_input, 'name'): # Check if it's an UploadedFile object
            input_type = "UploadedFile"
            file_name = getattr(image_input, 'name', 'N/A')
            # Use file_id and size for better tracking if available
            file_id = getattr(image_input, 'id', 'N/A')
            file_size = getattr(image_input, 'size', 'N/A')
            input_details = f"name='{file_name}', id='{file_id}', size={file_size}"
            logger.info(f"Processing uploaded file {input_details} for OCR.")

            # Ensure pointer is at the beginning BEFORE reading bytes
            try:
                image_input.seek(0)
                logger.debug(f"File pointer reset for {input_details}")
            except Exception as seek_err:
                logger.warning(f"Could not seek(0) on input file {input_details}: {seek_err}")
                # Continue anyway, getvalue() might still work

            image_bytes = image_input.getvalue() # Read fresh bytes

            if not image_bytes:
                logger.warning(f"extract_text_from_image received empty file bytes for {input_details}.")
                # Return empty text and None for image if bytes are empty
                return "", None

            try:
                # Open the image from the fresh bytes using BytesIO
                image = Image.open(io.BytesIO(image_bytes))
                logger.debug(f"Successfully opened image from bytes for {input_details}. Size: {image.size}, Mode: {image.mode}")
            except UnidentifiedImageError:
                st.error(f"Invalid or corrupted image file uploaded: {file_name}")
                logger.error(f"UnidentifiedImageError during Image.open for {input_details}.")
                return "", None
            except Exception as img_open_e:
                 st.error(f"Error opening image file {file_name}: {img_open_e}")
                 logger.error(f"Error during Image.open for {input_details}: {img_open_e}", exc_info=True)
                 return "", None
        else:
            input_type = str(type(image_input))
            logger.error(f"Unsupported input type for OCR: {input_type}")
            st.error("Internal error: Invalid data type received for image processing.")
            return "", None

        # Check if image was successfully loaded/passed
        if image is None:
             logger.error(f"Image object is None after input processing. Input type: {input_type}, Details: {input_details}")
             st.error("Internal error: Failed to load image data.")
             return "", None

        # --- Perform OCR ---
        try:
            # Ensure image is in RGB format for OpenCV compatibility
            pil_image_rgb = image.convert('RGB')
            img_cv = cv2.cvtColor(np.array(pil_image_rgb), cv2.COLOR_RGB2BGR)
            gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
            thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
            extracted_text = pytesseract.image_to_string(thresh, config='--psm 6')
            logger.info(f"OCR extracted {len(extracted_text)} characters from {input_details}.")
            # Return the originally opened/passed PIL image (not the cv2 one)
            return extracted_text, image
        except pytesseract.TesseractNotFoundError:
            st.error("Tesseract OCR engine not found. Please check the system configuration.")
            logger.critical("TesseractNotFoundError - Check system PATH or Tesseract installation in environment.")
            # Return the image even if OCR fails, it might be needed visually
            return "OCR Error: Tesseract not found", image
        except Exception as ocr_e:
            st.warning(f"An error occurred during text extraction (OCR) for {input_details}: {ocr_e}") # Use warning for non-critical OCR errors
            logger.warning(f"OCR Error processing {input_details}: {ocr_e}", exc_info=True)
            # Return the image even if OCR fails, with an error message as text
            return f"OCR Error: {ocr_e}", image

    except Exception as e:
        # Catch-all for unexpected errors during input handling phase
        st.error(f"An unexpected error occurred preparing the image for OCR: {e}")
        logger.error(f"Unexpected Error in extract_text_from_image: {e}", exc_info=True)
        return "", None

# REVISED ensure_fresh_control_image using session state for explicit caching control
def ensure_fresh_control_image(control_layout_file):
    """
    Ensures a fresh control image is processed only when the uploaded file changes.
    Uses st.session_state for explicit cache management based on file identity.
    """
    # Initialize session state keys if they don't exist
    if 'processed_control_image_key' not in st.session_state:
        st.session_state.processed_control_image_key = None
    if 'processed_control_image' not in st.session_state:
        st.session_state.processed_control_image = None

    if control_layout_file is None:
        # If no file is uploaded, ensure any previous state is cleared
        if st.session_state.processed_control_image_key is not None:
            st.session_state.processed_control_image_key = None
            st.session_state.processed_control_image = None
            logger.info("Control file removed. Cleared control image from session state.")
        return None

    # --- Generate a unique key for the current file ---
    # Best effort: use file_id if available, otherwise combine name and size.
    # Using file_id is generally more reliable than name+size.
    file_id = getattr(control_layout_file, 'id', None)
    file_name = getattr(control_layout_file, 'name', 'unknown')
    file_size = getattr(control_layout_file, 'size', 'unknown')

    if file_id:
        current_file_key = f"control_{file_id}"
    else:
        # Fallback if file_id is not available
        current_file_key = f"control_{file_name}_{file_size}"

    logger.debug(f"Ensuring control image. Current file key: '{current_file_key}', Stored key: '{st.session_state.processed_control_image_key}'")

    # --- Check session state ---
    # Check if the key matches the last processed key AND the image exists in state
    if (st.session_state.processed_control_image_key == current_file_key and
        isinstance(st.session_state.processed_control_image, Image.Image)):

        logger.info(f"Returning cached control image from session state for key: '{current_file_key}'")
        return st.session_state.processed_control_image
    else:
        # --- File has changed or not processed yet ---
        action = "Processing new" if st.session_state.processed_control_image_key is None else "Processing changed"
        logger.info(f"{action} control file. Key '{current_file_key}' requires processing.")

        # Call the robust extraction function
        _, control_image_input_pil = extract_text_from_image(control_layout_file)

        # Validate the result
        if not isinstance(control_image_input_pil, Image.Image):
            logger.error(f"Failed to process control layout image '{file_name}'. extract_text_from_image returned type: {type(control_image_input_pil)}")
            st.error(f"Could not load the control image: {file_name}. Please check the file or try uploading again.")
            # Clear potentially invalid state if processing failed
            st.session_state.processed_control_image_key = None
            st.session_state.processed_control_image = None
            return None
        else:
            # --- Store the newly processed image and its key in session state ---
            logger.info(f"Successfully processed fresh control image '{file_name}' -> PIL Image size: {control_image_input_pil.size}. Storing in session state with key '{current_file_key}'.")
            # Make a copy before storing in session state to prevent potential
            # downstream modifications affecting the cached version.
            st.session_state.processed_control_image = control_image_input_pil.copy()
            st.session_state.processed_control_image_key = current_file_key
            # Return the original (or copied) PIL image for immediate use
            return control_image_input_pil



# --- Price Extraction Function ---
# DO NOT CHANGE - Price extraction logic should remain dynamic
def extract_prices_from_test_type(test_type):
    """Extract control and variant prices from test description"""
    logger.info(f"Extracting prices from test description: {test_type}")

    # Pattern to find price pairs in different formats
    control_pattern = r'control:?\s*\$?(\d+\.?\d*)|control\s+\$?(\d+\.?\d*)'
    variant_pattern = r'variant:?\s*\$?(\d+\.?\d*)|variant\s+\$?(\d+\.?\d*)'

    # Simple $ pattern as fallback
    dollar_pattern = r'\$(\d+\.?\d*)'

    # Try to extract control price
    control_match = re.search(control_pattern, test_type, re.IGNORECASE)
    control_price = None
    if control_match:
        # Get the first non-None group from the regex match
        control_price = next((g for g in control_match.groups() if g is not None), None)
        logger.info(f"Found explicit control price: {control_price}")

    # Try to extract variant price
    variant_match = re.search(variant_pattern, test_type, re.IGNORECASE)
    variant_price = None
    if variant_match:
        # Get the first non-None group from the regex match
        variant_price = next((g for g in variant_match.groups() if g is not None), None)
        logger.info(f"Found explicit variant price: {variant_price}")

    # If we couldn't find explicit control/variant prices, fall back to simple $ pattern
    if not control_price or not variant_price:
        logger.info("Explicit control/variant price not found, searching for $ values...")
        all_prices = re.findall(dollar_pattern, test_type)
        logger.info(f"Found price values: {all_prices}")
        if len(all_prices) >= 2:
            control_price = all_prices[0]
            variant_price = all_prices[1]
            logger.info(f"Using first two $ values: control={control_price}, variant={variant_price}")
        elif len(all_prices) == 1:
            control_price = all_prices[0]
            # If we only found one price, assume the variant is slightly different
            try:
                variant_price = str(float(control_price) * 0.8)  # 20% discount as default
                # Round to 2 decimal places
                variant_price = str(round(float(variant_price), 2))
                logger.info(f"Found only one price, assuming 20% discount for variant: control={control_price}, variant={variant_price}")
            except ValueError:
                # If conversion fails, use a default
                variant_price = "5.00"
                logger.warning(f"Could not calculate variant price from {control_price}, using default: {variant_price}")

    # Set defaults if no prices found
    if not control_price:
        control_price = "7.95"
        logger.info(f"No control price found, using default: {control_price}")

    if not variant_price:
        variant_price = "5.00"
        logger.info(f"No variant price found, using default: {variant_price}")

    # Format prices with $ sign if not already present
    control_price_str = f"${control_price}" if not str(control_price).startswith("$") else control_price
    variant_price_str = f"${variant_price}" if not str(variant_price).startswith("$") else variant_price

    logger.info(f"Final extracted prices: Control={control_price_str}, Variant={variant_price_str}")
    return control_price_str, variant_price_str


def extract_metrics_from_supporting_data(image_obj):
    """Extract key metrics from a PIL image object using OCR."""
    internal_default_metrics = { "conversion_rate": "N/A", "total_checkout": "N/A", "checkouts": "N/A", "orders": "N/A", "shipping_revenue": "N/A", "aov": "N/A" }
    if not isinstance(image_obj, Image.Image):
        logger.warning("extract_metrics needs PIL object. Received non-image or None.")
        return internal_default_metrics
    try:
        # Use the robust extract_text_from_image function
        extracted_text, _ = extract_text_from_image(image_obj)
        if not extracted_text or extracted_text.startswith("OCR Error"):
            logger.warning(f"No text or OCR error from supporting data image for metrics. OCR Result: '{extracted_text}'")
            return internal_default_metrics

        # Existing regex patterns - kept as is
        conv_match = re.search(r'(?:conversion|checkout\s*conversion)\s*(?:rate)?[:\s]*(\d{1,3}(?:,\d{3})*\.?\d*%?|\d+\.?\d*%?)', extracted_text, re.IGNORECASE)
        total_co_match = re.search(r'(?:%|percent)\s*total\s*checkout[:\s]*(\d{1,3}(?:,\d{3})*\.?\d*%?|\d+\.?\d*%?)', extracted_text, re.IGNORECASE)
        checkouts_match = re.search(r'(?:#|number\s+of)?\s*Checkouts[:\s]*(\d{1,3}(?:,\d{3})*)', extracted_text, re.IGNORECASE)
        orders_match = re.search(r'(?:#|number\s+of)?\s*Orders[:\s]*(\d{1,3}(?:,\d{3})*)', extracted_text, re.IGNORECASE)
        ship_rev_match = re.search(r'(?:avg|average)?\s*shipping\s*revenue[:\s]*(\$\s*\d{1,3}(?:,\d{3})*\.?\d*)', extracted_text, re.IGNORECASE)
        aov_match = re.search(r'(?:AOV|average\s*order\s*value)[:\s]*(\$\s*\d{1,3}(?:,\d{3})*\.?\d*)', extracted_text, re.IGNORECASE)

        metrics = {
            "conversion_rate": conv_match.group(1).strip() if conv_match else internal_default_metrics["conversion_rate"],
            "total_checkout": total_co_match.group(1).strip() if total_co_match else internal_default_metrics["total_checkout"],
            "checkouts": checkouts_match.group(1).strip() if checkouts_match else internal_default_metrics["checkouts"],
            "orders": orders_match.group(1).strip() if orders_match else internal_default_metrics["orders"],
            "shipping_revenue": ship_rev_match.group(1).strip() if ship_rev_match else internal_default_metrics["shipping_revenue"],
            "aov": aov_match.group(1).strip() if aov_match else internal_default_metrics["aov"],
        }
        if metrics["conversion_rate"] != "N/A" and '%' not in metrics["conversion_rate"]: metrics["conversion_rate"] += "%"
        if metrics["total_checkout"] != "N/A" and '%' not in metrics["total_checkout"]: metrics["total_checkout"] += "%"

        logger.info(f"Extracted Metrics: {metrics}")
        return metrics
    except Exception as e:
        logger.error(f"Metric extraction error: {e}", exc_info=True)
        return internal_default_metrics

# --- HTML Variant Generation ---
# Kept as is, relates to variant generation, not control image handling
def generate_shipping_html(standard_price="$7.95", rush_price="$24.95", is_variant=False):
    """ Generate HTML content for shipping options display """
    # NOTE: The variant label in the HTML itself is kept for the image generation
    # but the corresponding label in the PPTX will be removed later.

    # The changes go INSIDE this multi-line f-string below:
    html = f"""<!DOCTYPE html>
<html>
<head>
<style>
    body {{ font-family: Arial, sans-serif; background-color: #f8f9fa; margin: 0; padding: 20px; box-sizing: border-box; }}
    .container {{ max-width: 580px; background-color: {PDQ_COLORS['white']}; border: 1px solid {PDQ_COLORS['html_border']}; border-radius: 6px; padding: 20px; position: relative; box-shadow: 0 1px 2px rgba(0,0,0,0.05); }}
    h2 {{ margin-top: 0; margin-bottom: 15px; font-size: 16px; font-weight: 600; color: {PDQ_COLORS['black']}; }}
    .shipping-option {{ border: 1px solid {PDQ_COLORS['html_border']}; border-radius: 6px; padding: 15px; margin-bottom: 10px; display: flex; align-items: flex-start; transition: all 0.2s ease-in-out; }}
    .shipping-option.selected {{ border-color: {PDQ_COLORS['html_selected_border']}; background-color: {PDQ_COLORS['html_selected_bg']}; }}
    .radio {{ margin-right: 12px; margin-top: 3px; flex-shrink: 0; }}
    .radio-dot {{ width: 18px; height: 18px; border-radius: 50%; display: flex; align-items: center; justify-content: center; }}
    .radio-selected .radio-dot {{ background-color: {PDQ_COLORS['html_selected_border']}; }}
    .radio-selected .radio-dot-inner {{ width: 7px; height: 7px; border-radius: 50%; background-color: white; }}
    .radio-unselected .radio-dot {{ border: 2px solid {PDQ_COLORS['html_radio_border']}; background-color: white; }}
    .shipping-details {{ flex-grow: 1; }}
    .shipping-title {{ font-weight: 600; font-size: 14px; margin-bottom: 4px; color: #333; }}
    .shipping-subtitle {{ color: {PDQ_COLORS['grey_text']}; font-size: 12px; }}

    /* --- MODIFY FONT SIZE HERE --- */
    .shipping-price {{
        font-weight: 600;
        /* font-size: 14px; */  /* Comment out or remove the original */
        font-size: 15px;      /* << NEW FONT SIZE (Adjust 15px/16px as needed) */
        text-align: right;
        min-width: 60px;
        color: #333;
        margin-left: 10px;
    }}

    .footnote {{ font-size: 12px; color: {PDQ_COLORS['grey_text']}; margin-top: 15px; }}

    /* --- AND MODIFY FONT SIZE HERE --- */
    .variant-label {{
        position: absolute;
        top: 10px;
        right: 10px;
        background-color: {PDQ_COLORS['white']};
        border: 1px solid {PDQ_COLORS['electric_violet']};
        color: {PDQ_COLORS['electric_violet']};
        font-weight: 600;
        /* font-size: 9px; */   /* Comment out or remove the original */
        font-size: 10px;       /* << NEW FONT SIZE (Adjust 10px/11px as needed) */
        padding: 2px 5px;
        border-radius: 3px;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }}
</style>
</head>
<body>
    <div class="container">
        <h2>Shipping method</h2>
        {f'<div class="variant-label">VARIANT</div>' if is_variant else ''}
        <div class="shipping-option selected">
            <div class="radio radio-selected"><div class="radio-dot"><div class="radio-dot-inner"></div></div></div>
            <div class="shipping-details">
                <div class="shipping-title">Standard Shipping & Processing* (4-7 Business Days)</div>
                <div class="shipping-subtitle">Please allow 1-2 business days for order processing</div>
            </div>
            <div class="shipping-price">{standard_price}</div>
        </div>
        <div class="shipping-option">
            <div class="radio radio-unselected"><div class="radio-dot"></div></div>
            <div class="shipping-details">
                <div class="shipping-title">Rush Shipping* (2 Business Days)</div>
                <div class="shipping-subtitle">Please allow 1-2 business days for order processing</div>
            </div>
            <div class="shipping-price">{rush_price}</div>
        </div>
        <div class="footnote">*Includes $1.49 processing fee</div>
    </div>
</body>
</html>"""
    return html

# Kept as is
def html_to_image(html_content, output_path="temp_shipping_image.png", size=(600, 350)): # Increased height
    """ Convert HTML content to an image using html2image """
    try:
        temp_dir = tempfile.gettempdir()
        # Add --headless=new flag as suggested by Chrome error logs
        hti = Html2Image(output_path=temp_dir, size=size, custom_flags=['--headless=new', '--no-sandbox', '--disable-gpu'])
        # Create a more unique filename to avoid potential collisions in temp dir
        unique_filename = f"{os.path.splitext(os.path.basename(output_path))[0]}_{datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')}.png"
        full_output_path = os.path.join(temp_dir, unique_filename)

        logger.info(f"Attempting HTML screenshot to: {full_output_path} with size {size}")
        paths = hti.screenshot(html_str=html_content, save_as=unique_filename)

        if not paths or not os.path.exists(paths[0]):
             raise RuntimeError(f"html2image failed to create the screenshot file at {paths[0] if paths else 'unknown path'}")

        logger.info(f"Screenshot successful: {paths[0]}")
        # Open, copy, then close immediately to release file handle
        with Image.open(paths[0]) as img:
            img_copy = img.copy()

        # Attempt to remove the temporary file
        try:
            os.remove(paths[0])
            logger.info(f"Removed temporary screenshot: {paths[0]}")
        except Exception as e:
            logger.warning(f"Failed to remove temp screenshot {paths[0]}: {e}")

        return img_copy

    except Exception as e:
        st.error(f"Error converting HTML to image using html2image: {e}")
        logger.error(f"html_to_image error: {e}", exc_info=True)
        # Generate placeholder on error
        placeholder = Image.new('RGB', size, color=hex_to_rgb(PDQ_COLORS["moon_raker"]))
        draw = ImageDraw.Draw(placeholder)
        try:
            font = get_font("Arial", 14)
            draw.text((10, 10), "Error generating image preview", fill=(0,0,0), font=font)
        except Exception: pass # Ignore font errors for placeholder
        return placeholder

def extract_prices_from_test_type(test_type):
    """Extract control and variant prices from test description"""
    logger.info(f"Extracting prices from test description: {test_type}")
    
    # Pattern to find price pairs in different formats
    control_pattern = r'control:?\s*\$?(\d+\.?\d*)|control\s+\$?(\d+\.?\d*)'
    variant_pattern = r'variant:?\s*\$?(\d+\.?\d*)|variant\s+\$?(\d+\.?\d*)'
    
    # Simple $ pattern as fallback
    dollar_pattern = r'\$(\d+\.?\d*)'
    
    # Try to extract control price
    control_match = re.search(control_pattern, test_type, re.IGNORECASE)
    control_price = None
    if control_match:
        # Get the first non-None group from the regex match
        control_price = next((g for g in control_match.groups() if g is not None), None)
        logger.info(f"Found explicit control price: {control_price}")
    
    # Try to extract variant price
    variant_match = re.search(variant_pattern, test_type, re.IGNORECASE)
    variant_price = None
    if variant_match:
        # Get the first non-None group from the regex match
        variant_price = next((g for g in variant_match.groups() if g is not None), None)
        logger.info(f"Found explicit variant price: {variant_price}")
    
    # If we couldn't find explicit control/variant prices, fall back to simple $ pattern
    if not control_price or not variant_price:
        logger.info("Explicit control/variant price not found, searching for $ values...")
        all_prices = re.findall(dollar_pattern, test_type)
        logger.info(f"Found price values: {all_prices}")
        if len(all_prices) >= 2:
            control_price = all_prices[0]
            variant_price = all_prices[1]
            logger.info(f"Using first two $ values: control={control_price}, variant={variant_price}")
        elif len(all_prices) == 1:
            control_price = all_prices[0]
            # If we only found one price, assume the variant is slightly different
            try:
                variant_price = str(float(control_price) * 0.8)  # 20% discount as default
                # Round to 2 decimal places
                variant_price = str(round(float(variant_price), 2))
                logger.info(f"Found only one price, assuming 20% discount for variant: control={control_price}, variant={variant_price}")
            except ValueError:
                # If conversion fails, use a default
                variant_price = "5.00"
                logger.warning(f"Could not calculate variant price from {control_price}, using default: {variant_price}")
    
    # Set defaults if no prices found
    if not control_price:
        control_price = "7.95"
        logger.info(f"No control price found, using default: {control_price}")
    
    if not variant_price:
        variant_price = "5.00"
        logger.info(f"No variant price found, using default: {variant_price}")
    
    # Format prices with $ sign if not already present
    control_price_str = f"${control_price}" if not str(control_price).startswith("$") else control_price
    variant_price_str = f"${variant_price}" if not str(variant_price).startswith("$") else variant_price
    
    logger.info(f"Final extracted prices: Control={control_price_str}, Variant={variant_price_str}")
    return control_price_str, variant_price_str
# Kept as is
def generate_shipping_options(old_price="$7.95", new_price="$5.00"):
    """ Generate control and variant shipping option images """
    logger.info(f"Generating shipping HTML for control and variant with prices: control={old_price}, variant={new_price}")
    # Create a unique timestamp for this generation
    timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')

    # Generate unique filenames with timestamp
    control_filename = f"control_shipping_{timestamp}.png"
    variant_filename = f"variant_shipping_{timestamp}.png"

    control_html = generate_shipping_html(old_price, "$24.95", is_variant=False)
    variant_html = generate_shipping_html(new_price, "$24.95", is_variant=True)
    logger.info("Converting HTML to images...")

    # Use the timestamped filenames
    control_image = html_to_image(control_html, output_path=control_filename)
    variant_image = html_to_image(variant_html, output_path=variant_filename)

    # Log the image sizes after generation
    logger.info(f"Generated control shipping image with price {old_price} and size: {control_image.size if control_image else 'None'}")
    logger.info(f"Generated variant shipping image with price {new_price} and size: {variant_image.size if variant_image else 'None'}")

    logger.info("Shipping option image generation complete.")
    return control_image, variant_image

# Kept as is
def create_price_variant(old_price, new_price):
    """
    Create control and variant images based on the HTML template with different prices.
    This function ensures we use the HTML-based generation for a consistent look.

    Args:
        old_price: Original price string (e.g. "$7.95") - Used for the control image
        new_price: New price string (e.g. "$5.00") - Used for the variant image

    Returns:
        tuple: (control_image, variant_image) as PIL Image objects
    """
    logger.info(f"Creating HTML-based variant image by changing price from {old_price} (Control) to {new_price} (Variant)")
    try:
        # Use the generate_shipping_options function which calls html_to_image
        # Note: generate_shipping_options as provided earlier seems to use the *first* price for standard shipping
        # in *both* control and variant HTML, and the *second* price for rush shipping.
        # Reverting to the interpretation from the original code which mapped old_price to control standard, new_price to variant standard.
        def generate_shipping_options_html(standard_price_control, standard_price_variant, rush_price="$24.95"):
            logger.info(f"Generating shipping HTML for control and variant with prices: control_std={standard_price_control}, variant_std={standard_price_variant}, rush={rush_price}")
            timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')
            control_filename = f"control_shipping_{timestamp}.png"
            variant_filename = f"variant_shipping_{timestamp}.png"

            # Control HTML: standard shipping is the 'old_price'
            control_html = generate_shipping_html(standard_price_control, rush_price, is_variant=False)
            # Variant HTML: standard shipping is the 'new_price'
            variant_html = generate_shipping_html(standard_price_variant, rush_price, is_variant=True)

            logger.info("Converting HTML to images...")
            control_image = html_to_image(control_html, output_path=control_filename)
            variant_image = html_to_image(variant_html, output_path=variant_filename)

            logger.info(f"Generated control shipping image with price {standard_price_control} and size: {control_image.size if control_image else 'None'}")
            logger.info(f"Generated variant shipping image with price {standard_price_variant} and size: {variant_image.size if variant_image else 'None'}")
            logger.info("Shipping option image generation complete.")
            return control_image, variant_image

        # Call the internal HTML generation using the prices
        control_image, variant_image = generate_shipping_options_html(old_price, new_price)


        if not isinstance(control_image, Image.Image) or not isinstance(variant_image, Image.Image):
            raise ValueError("Failed to generate valid images from HTML")

        logger.info(f"create_price_variant successful. Control size: {control_image.size}, Variant size: {variant_image.size}")
        return control_image, variant_image

    except Exception as e:
        logger.error(f"Error in create_price_variant: {e}", exc_info=True)
        st.error("An error occurred while generating the price variant images.")
        # In case of error, return placeholder images
        size = (600, 350)
        placeholder = Image.new('RGB', size, color=hex_to_rgb(PDQ_COLORS["moon_raker"]))
        draw = ImageDraw.Draw(placeholder)
        try:
            font = get_font("Arial", 14)
            draw.text((10, 10), f"Error generating price image.", fill=(0,0,0), font=font)
            draw.text((10, 30), f"Control Price: {old_price}", fill=(0,0,0), font=font)
            draw.text((10, 50), f"Variant Price: {new_price}", fill=(0,0,0), font=font)
        except:
            draw.text((10, 10), f"Error generating price image.", fill=(0,0,0))

        control_copy = placeholder.copy()
        variant_copy = placeholder.copy()
        draw_v = ImageDraw.Draw(variant_copy)
        try:
            font_bold = get_font("Arial", 16, bold=True)
            text_width, text_height = draw_v.textbbox((0,0), "VARIANT", font=font_bold)[2:]
            draw_v.text((size[0] - text_width - 20, 20), "VARIANT", fill=(0,0,0), font=font_bold)
        except:
            draw_v.text((size[0] - 80, 20), "VARIANT", fill=(0,0,0))

        return control_copy, variant_copy

# --- PDF Processing ---
# Kept as is
def extract_from_pdf(pdf_file):
    """Extract content from PDF files using PyMuPDF"""
    pdf_content = []
    if not pdf_file:
        logger.warning("extract_from_pdf received None.")
        return pdf_content
    try:
        logger.info(f"Processing PDF: {getattr(pdf_file, 'name', 'N/A')}")
        # Ensure pointer is at beginning before reading
        pdf_file.seek(0)
        pdf_bytes = pdf_file.read()
        if not pdf_bytes:
            logger.warning("Empty PDF file uploaded.")
            return pdf_content

        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            logger.info(f"Opened PDF with {len(doc)} pages.")
            for page_num, page in enumerate(doc):
                page_text = page.get_text("text")
                page_images = []
                image_list = page.get_images(full=True)
                logger.debug(f"Page {page_num+1}: Found {len(image_list)} raw image entries.")
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    if xref == 0: continue # Skip if xref is 0 (indicates invalid image entry)
                    try:
                        base_image = doc.extract_image(xref)
                        if not base_image:
                            logger.warning(f"extract_image returned None for xref {xref} on page {page_num+1}.")
                            continue
                        image_bytes = base_image["image"]
                        if image_bytes:
                            try:
                                # Open image from bytes
                                image = Image.open(io.BytesIO(image_bytes))
                                page_images.append(image)
                                logger.debug(f"Successfully extracted image {img_index+1} (xref: {xref}) from page {page_num+1}.")
                            except UnidentifiedImageError:
                                logger.warning(f"Skipping unidentified image {img_index+1} (xref: {xref}) on page {page_num+1}.")
                            except Exception as img_open_e:
                                logger.warning(f"Error opening image {img_index+1} (xref: {xref}) from page {page_num+1}: {img_open_e}")
                        else:
                            logger.warning(f"Empty image bytes for xref {xref} on page {page_num+1}.")
                    except Exception as e:
                        # Catch errors during extraction for a specific image
                        logger.warning(f"PDF Image Extraction Warning - Page {page_num+1}, Img Index {img_index+1} (xref: {xref}): {e}")

                pdf_content.append({"page_num": page_num + 1, "text": page_text, "images": page_images})
                logger.debug(f"Processed Page {page_num+1}: Text length {len(page_text)}, Images found: {len(page_images)}")

        logger.info(f"Finished processing PDF: Found {sum(len(p['images']) for p in pdf_content)} valid images total.")
    except Exception as e:
        st.error(f"PDF Processing Error: {e}")
        logger.error(f"PDF Error: {e}", exc_info=True)
    return pdf_content

# --- Slide Content Generation Helper ---
# Kept as is
class PDQSlideGeneratorHelper:
    """ Helper class for generating slide content like hypothesis, KPIs, etc. """
    def __init__(self):
        """Initialize the helper"""
        self.hypothesis_templates = {
             "price": "We believe that adjusting pricing for {segment} will {expected_outcome} because {rationale}.",
             "shipping": "We believe that modifying shipping options for {segment} will {expected_outcome} because {rationale}.",
             "layout": "We believe that updating the interface design for {segment} will {expected_outcome} because {rationale}.",
             "messaging": "We believe that changing the messaging for {segment} will {expected_outcome} because {rationale}.",
             "generic": "We believe that the proposed changes for {segment} will {expected_outcome} because {rationale}."
        }
        self.outcome_templates = {
             "price": ["increase revenue without significantly impacting conversion rate", "optimize revenue per visitor", "improve average order value"],
             "shipping": ["improve conversion rates by addressing shipping concerns", "reduce cart abandonment related to shipping", "increase customer satisfaction with delivery options"],
             "layout": ["improve user engagement and streamline the checkout flow", "reduce friction and increase task completion rate", "enhance visual appeal and clarity"],
             "messaging": ["increase click-through rates on key actions", "improve conversion by clarifying value proposition", "enhance perceived trust and urgency"],
             "generic": ["positively impact key performance indicators", "improve user experience and drive conversions", "optimize the user journey"]
        }
        self.rationale_templates = {
             "price": ["analysis shows potential for price optimization", "competitor pricing suggests room for adjustment", "segment behavior indicates willingness to pay"],
             "shipping": ["customer feedback indicates shipping costs are a pain point", "data suggests high drop-off at the shipping stage", "offering faster options may appeal to this segment"],
             "layout": ["the current layout has known usability issues", "best practices suggest the new design will perform better", "heatmaps indicate users struggle with the current interface"],
             "messaging": ["the new copy is clearer and more benefit-oriented", "current messaging is underperforming in preliminary tests", "aligning messaging with segment needs should improve resonance"],
             "generic": ["this change addresses a key hypothesis from user research", "market trends support this type of modification", "data indicates an opportunity for improvement in this area"]
        }
    def generate_hypothesis(self, test_type, segment, supporting_data_text=""):
        category = "generic"; test_type_lower = test_type.lower()
        if any(word in test_type_lower for word in ["price", "$", "cost", "charge", "fee"]): category = "price"
        elif any(word in test_type_lower for word in ["shipping", "delivery", "freight"]): category = "shipping"
        elif any(word in test_type_lower for word in ["layout", "design", "ui", "ux", "interface", "position"]): category = "layout"
        elif any(word in test_type_lower for word in ["message", "copy", "text", "wording", "cta", "button", "headline"]): category = "messaging"
        expected_outcome = random.choice(self.outcome_templates[category]); rationale = random.choice(self.rationale_templates[category])
        hypothesis = self.hypothesis_templates[category].format( segment=segment if segment else "users", expected_outcome=expected_outcome, rationale=rationale )
        logger.info(f"Generated Hypothesis (Category: {category}): {hypothesis}")
        return hypothesis
    def parse_test_type(self, test_type):
        parts = test_type.split('—'); title = test_type.strip()
        if len(parts) > 1:
             # If there's an em dash, take the part before it as the main title
             title = parts[0].strip()
             # Potentially use parts[1] as subtitle/description if needed later
             logger.info(f"Parsed test type: Title='{title}', Details='{parts[1].strip()}'")
        else:
             logger.info(f"Parsed test type (no dash): Title='{title}'")
        return title

    def infer_goals_and_kpis(self, test_type):
        test_type_lower = test_type.lower(); goal, kpi, impact = "Improve Performance", "Conversion Rate (CVR)", "3-5%"
        if any(term in test_type_lower for term in ["price", "pricing", "cost", "$", "revenue", "aov", "value"]): goal, kpi, impact = "Increase Revenue", "Revenue Per Visitor (RPV)", "5-8%"
        elif any(term in test_type_lower for term in ["conversion", "cvr", "checkout", "completion", "purchase"]): goal, kpi, impact = "Increase Conversion Rate", "Conversion Rate (CVR)", "3-5%"
        elif any(term in test_type_lower for term in ["shipping", "delivery", "ship", "abandonment"]): goal, kpi, impact = "Optimize Shipping", "Checkout Completion Rate", "4-7%"
        elif any(term in test_type_lower for term in ["layout", "design", "ui", "interface", "ux", "engagement"]): goal, kpi, impact = "Improve User Experience", "Engagement Rate / CVR", "5-10%"
        elif any(term in test_type_lower for term in ["message", "copy", "text", "wording", "cta", "click", "ctr"]): goal, kpi, impact = "Increase Engagement", "Click-Through Rate (CTR)", "8-15%"
        logger.info(f"Inferred Goal: {goal}, KPI: {kpi}, Impact: {impact}")
        return goal, kpi, impact
    def generate_tags(self, test_type, segment, supporting_data_text=""):
        tags = set(); combined_text = f"{test_type} {segment} {supporting_data_text}".lower()
        tag_map = {
            "Price Sensitivity": ["price", "$", "cost", "fee", "charge"], "Shipping Options": ["shipping", "delivery", "ship", "freight"],
            "UI/UX Design": ["layout", "design", "ui", "ux", "interface", "visual"], "Messaging/Copy": ["message", "copy", "text", "wording", "cta", "headline", "button"],
            "Checkout Process": ["checkout", "cart", "payment", "completion"], "Revenue Optimization": ["revenue", "aov", "rpv", "value"],
            "Conversion Rate Optimization": ["conversion", "cvr", "purchase"], "Mobile Experience": ["mobile", "smartphone", "ios", "android"],
            "Desktop Experience": ["desktop", "pc"], "New Customers": ["first time", "new user", "new customer", "acquisition"],
            "Returning Customers": ["returning", "repeat", "loyal", "retention"], "Cart Abandonment": ["abandoned", "abandonment", "drop-off"],
            "Free Shipping Threshold": ["fst", "free shipping"], "Urgency/Scarcity": ["urgent", "limited", "timer", "stock"], "Social Proof": ["review", "testimonial", "rating"],
        }
        for tag, keywords in tag_map.items():
            if any(keyword in combined_text for keyword in keywords): tags.add(tag)
        final_tags = list(tags)
        logger.info(f"Generated Tags: {final_tags[:4]}")
        return final_tags[:4]
    def determine_success_criteria(self, test_type, kpi, goal):
         criteria = f"Statistically significant improvement in {kpi}"
         if "revenue" in goal.lower() or "aov" in goal.lower(): criteria = f"Increase in {kpi} without significant negative impact on CVR"
         elif "conversion" in goal.lower(): criteria = f"Uplift of 1-3% in {kpi} with 85%+ confidence"
         elif "shipping" in goal.lower() or "abandonment" in goal.lower(): criteria = f"Decrease in cart abandonment rate or increase in {kpi}"
         elif "engagement" in goal.lower() or "experience" in goal.lower(): criteria = f"Improvement in engagement metrics (e.g., {kpi})"
         logger.info(f"Determined Success Criteria: {criteria}")
         return criteria

# --- NEW FUNCTION: Ensure fresh control image processing ---
def ensure_fresh_control_image(control_layout_file):
    """Ensures that we always process a fresh control image without caching issues."""
    if control_layout_file is None:
        logger.warning("No control layout file provided.")
        return None
    
    logger.info(f"Processing control layout file: {getattr(control_layout_file, 'name', 'unknown')}")
    _, control_image_input_pil = extract_text_from_image(control_layout_file)
    
    if not isinstance(control_image_input_pil, Image.Image):
        logger.error("Failed to process control layout image.")
        return None
        
    logger.info(f"Successfully processed fresh control image with size: {control_image_input_pil.size if control_image_input_pil else 'None'}")
    return control_image_input_pil

# --- PowerPoint Generation Function ---
def create_proper_pptx(title, hypothesis, segment, goal, kpi_impact_str, elements_tags,
                       timeline_str, success_criteria, checkouts_required_str,
                       control_image, variant_image, supporting_data_image=None):
    """Creates a PowerPoint slide precisely matching the reference image layout."""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
        slide_layout = prs.slide_layouts[6]; slide = prs.slides.add_slide(slide_layout)
        background = slide.background; fill = background.fill; fill.solid(); fill.fore_color.rgb = hex_to_rgbcolor(PDQ_COLORS["revolver"])

        # --- Layout Constants (Adjusted for taller images) ---
        LEFT_MARGIN = Inches(0.4); TOP_MARGIN = Inches(0.4); RIGHT_MARGIN = Inches(0.4); BOTTOM_MARGIN = Inches(0.3)
        GRID_BOX_WIDTH = Inches(2.0)
        GRID_BOX_HEIGHT = Inches(1.5)
        GRID_GAP_HORZ = Inches(0.18)
        GRID_GAP_VERT = Inches(0.2)
        HYPOTHESIS_TOP = Inches(1.3); HYPOTHESIS_HEIGHT = Inches(1.5)
        GRID_TOP = HYPOTHESIS_TOP + HYPOTHESIS_HEIGHT + GRID_GAP_VERT * 1.5
        TOTAL_GRID_WIDTH = (GRID_BOX_WIDTH * 4) + (GRID_GAP_HORZ * 3)

        RIGHT_COL_LEFT = LEFT_MARGIN + TOTAL_GRID_WIDTH + GRID_GAP_HORZ
        RIGHT_COL_WIDTH = prs.slide_width - RIGHT_COL_LEFT - RIGHT_MARGIN - Inches(0.2)

        hyp_width = TOTAL_GRID_WIDTH
        CONTROL_VAR_TITLE_HEIGHT=Inches(0.3); CONTROL_VAR_GAP=Inches(0.1)

        # --- Vertical positioning (Adjusted for taller images) ---
        CONTROL_TITLE_TOP = HYPOTHESIS_TOP
        CONTROL_CONTAINER_HEIGHT = Inches(2.1) # Increased height
        CONTROL_CONTAINER_TOP = CONTROL_TITLE_TOP + CONTROL_VAR_TITLE_HEIGHT + CONTROL_VAR_GAP

        VARIANT_CONTAINER_HEIGHT = Inches(2.1) # Increased height
        # Adjust variant top based on new control height + gap
        VARIANT_TITLE_TOP = CONTROL_CONTAINER_TOP + CONTROL_CONTAINER_HEIGHT + Inches(0.2)
        VARIANT_CONTAINER_TOP = VARIANT_TITLE_TOP + CONTROL_VAR_TITLE_HEIGHT + CONTROL_VAR_GAP

        logo_size=Inches(0.6)

        # --- Helper: Add Rounded Rect Shape with Text ---
        def add_rounded_rect_with_text(left, top, width, height, bg_color_hex, title_text="", title_icon="", title_color_hex=PDQ_COLORS["melrose"], title_font_size=Pt(12.5), content_text="", content_color_hex=PDQ_COLORS["magnolia"], content_font_size=Pt(10.5), content_align=PP_ALIGN.LEFT, title_align=PP_ALIGN.LEFT, bold_title=True, center_content_vertical=False):
            try:
                # Ensure dimensions are integers for pptx
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(width), int(height))
                shape.fill.solid(); shape.fill.fore_color.rgb = hex_to_rgbcolor(bg_color_hex); shape.line.fill.background()
                tf = shape.text_frame; tf.word_wrap = True;
                tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15);
                tf.margin_top = 0; tf.margin_bottom = 0
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

                tf.clear()
                p_title = tf.add_paragraph(); p_title.text = f"{title_icon} {title_text}".strip();
                p_title.font.name = 'Segoe UI'; p_title.font.size = title_font_size;
                p_title.font.bold = bold_title; p_title.font.color.rgb = hex_to_rgbcolor(title_color_hex);
                p_title.alignment = title_align; p_title.space_after = 0; p_title.space_before = 0

                if content_text:
                    p_title.line_spacing = 0.8
                    lines = content_text.split('\n');
                    first_para = tf.add_paragraph()
                    first_para.alignment = content_align
                    first_para.text = lines[0] if lines else ""
                    first_para.font.name = 'Segoe UI'
                    first_para.font.size = content_font_size
                    first_para.font.color.rgb = hex_to_rgbcolor(content_color_hex)
                    first_para.space_before = 0
                    first_para.space_after = 0
                    first_para.line_spacing = 0.9

                    for i, line in enumerate(lines[1:], 1):
                        current_para = tf.add_paragraph()
                        current_para.alignment = content_align
                        current_para.text = line
                        current_para.font.name = 'Segoe UI'
                        current_para.font.size = content_font_size
                        current_para.font.color.rgb = hex_to_rgbcolor(content_color_hex)
                        current_para.space_before = 0
                        current_para.space_after = 0
                        current_para.line_spacing = 0.9

                if center_content_vertical:
                    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    p_title.space_before = 0; p_title.space_after = 0
                    if len(tf.paragraphs)>1: tf.paragraphs[1].space_before = 0

                return shape
            except Exception as e: logger.error(f"Error adding rounded rect '{title_text}': {e}", exc_info=True); return None

        # --- Helper: Place Image Inside a Shape ---
        def place_image_in_shape(img_obj, target_shape, slide_shapes, context="image", padding=Inches(0.1), scale_factor=1.0):
            """Places and resizes a PIL image inside a target pptx shape."""
            if not isinstance(img_obj, Image.Image): logger.warning(f"Invalid image type for {context}"); return None
            if target_shape is None: logger.warning(f"Target shape is None for {context}"); return None

            try:
                if isinstance(padding, (list, tuple)) and len(padding) == 4: pad_l, pad_t, pad_r, pad_b = padding
                else: pad_l = pad_t = pad_r = pad_b = padding

                inner_left = target_shape.left + pad_l; inner_top = target_shape.top + pad_t
                inner_width = (target_shape.width - pad_l - pad_r) * scale_factor
                inner_height = (target_shape.height - pad_t - pad_b) * scale_factor

                if inner_width <= 0 or inner_height <= 0: logger.error(f"Invalid inner bounds for {context} in shape: W={inner_width}, H={inner_height}"); return None

                logger.info(f"Placing {context}: Original size {img_obj.size}. Target Bounds: L={inner_left/Inches(1):.2f}\", T={inner_top/Inches(1):.2f}\", W={inner_width/Inches(1):.2f}\", H={inner_height/Inches(1):.2f}\"")
                img_byte_arr = io.BytesIO(); img_obj.save(img_byte_arr, format='PNG'); img_byte_arr = img_byte_arr.getvalue()
                if not img_byte_arr: raise ValueError("Image saving to bytes failed.")

                try:
                    # Add picture with initial width, height will be calculated
                    pic = slide_shapes.add_picture(io.BytesIO(img_byte_arr), int(inner_left), int(inner_top), width=int(inner_width)) # Cast initial position/width
                except Exception as add_pic_err:
                    logger.error(f"add_picture failed for {context}: {add_pic_err}", exc_info=True)
                    return None # Stop if adding picture fails

                # Calculate aspect ratio AFTER adding the picture (it might adjust width)
                img_ratio = pic.height / pic.width if pic.width > 0 else 1 # Use actual pic dimensions

                # Adjust height based on width and ratio
                pic.height = int(pic.width * img_ratio)

                # Center the image - CAST TO INT HERE
                pic.left = int(inner_left + ((target_shape.width - pad_l - pad_r) - pic.width) / 2)
                pic.top = int(inner_top + ((target_shape.height - pad_t - pad_b) - pic.height) / 2) if inner_height > pic.height else int(inner_top)

                logger.info(f"Successfully placed/resized {context}: Final Size W={pic.width/Inches(1):.2f}\", H={pic.height/Inches(1):.2f}\" at L={pic.left/Inches(1):.2f}\", T={pic.top/Inches(1):.2f}\"")
                return pic
            except Exception as e:
                logger.error(f"Error placing {context} image in shape: {e}", exc_info=True)
                raise e

        # --- Build Slide Elements ---
        logo_box = add_rounded_rect_with_text(LEFT_MARGIN, TOP_MARGIN, logo_size, logo_size, PDQ_COLORS["electric_violet"])

        title_left = LEFT_MARGIN + logo_size + Inches(0.2);
        title_width = hyp_width - (logo_size + Inches(0.2));
        title_box = slide.shapes.add_textbox(int(title_left), int(TOP_MARGIN), int(title_width), int(logo_size)); # Cast dimensions
        tf_title = title_box.text_frame;
        tf_title.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE;
        p_title = tf_title.add_paragraph();
        p_title.text = title;
        p_title.font.size = Pt(30);
        p_title.font.bold = True; p_title.font.name = 'Segoe UI';
        p_title.font.color.rgb = hex_to_rgbcolor(PDQ_COLORS["white"])

        add_rounded_rect_with_text(LEFT_MARGIN, HYPOTHESIS_TOP, hyp_width, HYPOTHESIS_HEIGHT,
                                  PDQ_COLORS["valentino"], title_text="Hypothesis",
                                  title_icon="✓", content_text=hypothesis, content_font_size=Pt(11))

        grid_col_starts = [LEFT_MARGIN + i * (GRID_BOX_WIDTH + GRID_GAP_HORZ) for i in range(4)]
        grid_row2_top = GRID_TOP + GRID_BOX_HEIGHT + GRID_GAP_VERT

        add_rounded_rect_with_text(grid_col_starts[0], GRID_TOP, GRID_BOX_WIDTH, GRID_BOX_HEIGHT,
                                  PDQ_COLORS["valentino"], title_text="Segment", title_icon="👥", content_text=segment)
        add_rounded_rect_with_text(grid_col_starts[1], GRID_TOP, GRID_BOX_WIDTH, GRID_BOX_HEIGHT,
                                  PDQ_COLORS["valentino"], title_text="Timeline", title_icon="📅", content_text=timeline_str)
        add_rounded_rect_with_text(grid_col_starts[2], GRID_TOP, GRID_BOX_WIDTH, GRID_BOX_HEIGHT,
                                  PDQ_COLORS["valentino"], title_text="Goal", title_icon="🎯", content_text=goal)
        elements_text = "\n".join(elements_tags) if elements_tags else "N/A";
        add_rounded_rect_with_text(grid_col_starts[3], GRID_TOP, GRID_BOX_WIDTH, GRID_BOX_HEIGHT,
                                  PDQ_COLORS["valentino"], title_text="Elements", title_icon="", content_text=elements_text)

        support_data_width = (GRID_BOX_WIDTH * 2) + GRID_GAP_HORZ;
        support_data_height = GRID_BOX_HEIGHT * 1.5
        support_box_shape = add_rounded_rect_with_text(grid_col_starts[0], grid_row2_top, support_data_width,
                                                     support_data_height, PDQ_COLORS["valentino"],
                                                     title_text="Supporting Data", title_icon="✓")
        if supporting_data_image and support_box_shape:
            support_img_padding_horz = Inches(0.25)
            support_img_title_clearance = Inches(0.8)
            support_img_padding_bottom = Inches(0.25)
            place_image_in_shape(
                supporting_data_image, support_box_shape, slide.shapes, context="supporting_data",
                padding=(support_img_padding_horz, support_img_title_clearance,
                         support_img_padding_horz, support_img_padding_bottom),
                scale_factor=0.7
            )

        add_rounded_rect_with_text(grid_col_starts[2], grid_row2_top, GRID_BOX_WIDTH, GRID_BOX_HEIGHT,
                                  PDQ_COLORS["valentino"], title_text="Success Criteria",
                                  title_icon="✓", content_text=success_criteria)

        add_rounded_rect_with_text(grid_col_starts[3], grid_row2_top, GRID_BOX_WIDTH, GRID_BOX_HEIGHT,
                                  PDQ_COLORS["electric_violet"], title_text=f"🛍️ {checkouts_required_str}",
                                  title_icon="", title_font_size=Pt(20),
                                  title_color_hex=PDQ_COLORS["lightning_yellow"], title_align=PP_ALIGN.CENTER,
                                  content_text="# of checkouts\nrequired", content_font_size=Pt(9.5),
                                  content_align=PP_ALIGN.CENTER, center_content_vertical=True)

        img_padding = Inches(0.07)
        # Use the adjusted heights
        control_container_with_title_top = CONTROL_TITLE_TOP
        control_container_with_title_height = CONTROL_CONTAINER_HEIGHT + CONTROL_VAR_TITLE_HEIGHT + CONTROL_VAR_GAP
        control_container = add_rounded_rect_with_text(RIGHT_COL_LEFT, control_container_with_title_top,
                                                      RIGHT_COL_WIDTH, control_container_with_title_height,
                                                      PDQ_COLORS["white"]) # Use helper for consistency

        control_title_left = RIGHT_COL_LEFT + Inches(0.15)
        control_title_box = slide.shapes.add_textbox(int(control_title_left), int(CONTROL_TITLE_TOP - Inches(0.05)),
                                                   int(RIGHT_COL_WIDTH - Inches(0.3)), int(CONTROL_VAR_TITLE_HEIGHT)) # Cast dimensions
        ctrl_tf = control_title_box.text_frame
        ctrl_tf.margin_left=0; ctrl_tf.margin_right=0; ctrl_tf.margin_top=0; ctrl_tf.margin_bottom=0
        ctrl_p = ctrl_tf.add_paragraph()
        ctrl_p.text = "Control"
        ctrl_p.font.size = Pt(14)
        ctrl_p.font.bold = True; ctrl_p.font.name = 'Segoe UI'
        ctrl_p.font.color.rgb = hex_to_rgbcolor(PDQ_COLORS["black"])

        control_image_top = CONTROL_TITLE_TOP + CONTROL_VAR_TITLE_HEIGHT + CONTROL_VAR_GAP
        if control_container: # Check if container was created
            place_image_in_shape(control_image, control_container, slide.shapes,
                                context="control_shipping",
                                padding=(img_padding, control_image_top - control_container_with_title_top + img_padding,
                                        img_padding, img_padding),
                                scale_factor=0.95)

        # Use the adjusted heights and top position
        variant_container_with_title_top = VARIANT_TITLE_TOP
        variant_container_with_title_height = VARIANT_CONTAINER_HEIGHT + CONTROL_VAR_TITLE_HEIGHT + CONTROL_VAR_GAP
        variant_container = add_rounded_rect_with_text(RIGHT_COL_LEFT, variant_container_with_title_top,
                                                      RIGHT_COL_WIDTH, variant_container_with_title_height,
                                                      PDQ_COLORS["white"])

        variant_title_left = RIGHT_COL_LEFT + Inches(0.15)
        variant_title_box = slide.shapes.add_textbox(int(variant_title_left), int(VARIANT_TITLE_TOP - Inches(0.05)),
                                                   int(RIGHT_COL_WIDTH - Inches(0.3)), int(CONTROL_VAR_TITLE_HEIGHT)) # Cast dimensions
        var_tf = variant_title_box.text_frame
        var_tf.margin_left=0; var_tf.margin_right=0; var_tf.margin_top=0; var_tf.margin_bottom=0
        var_p = var_tf.add_paragraph()
        var_p.text = "Variant B (example)"
        var_p.font.size = Pt(14)
        var_p.font.bold = True; var_p.font.name = 'Segoe UI'
        var_p.font.color.rgb = hex_to_rgbcolor(PDQ_COLORS["black"])

        variant_image_top = VARIANT_TITLE_TOP + CONTROL_VAR_TITLE_HEIGHT + CONTROL_VAR_GAP
        if variant_container: # Check if container was created
            var_pic = place_image_in_shape(variant_image, variant_container, slide.shapes,
                                          context="variant_shipping",
                                          padding=(img_padding, variant_image_top - variant_container_with_title_top + img_padding,
                                                  img_padding, img_padding),
                                          scale_factor=0.95)

        # --- VARIANT TAG REMOVED ---

        # --- Footer with dynamic date ---
        footer_top = prs.slide_height - BOTTOM_MARGIN - Inches(0.55)
        footer_box = slide.shapes.add_textbox(int(LEFT_MARGIN), int(footer_top),
                                            int(prs.slide_width - LEFT_MARGIN - RIGHT_MARGIN), int(Inches(0.25)))
        footer_frame = footer_box.text_frame
        footer_frame.margin_bottom = 0
        footer_para = footer_frame.add_paragraph()

        # Use current date instead of hardcoded date
        current_date = datetime.datetime.now().strftime('%B %d, %Y')
        footer_para.text = f"PDQ A/B Test | {current_date} | Confidential"
        footer_para.font.size = Pt(9)
        footer_para.font.italic = True
        footer_para.font.name = 'Segoe UI'
        footer_para.alignment = PP_ALIGN.RIGHT
        footer_para.font.color.rgb = hex_to_rgbcolor(PDQ_COLORS["melrose"])

        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        logger.info("PPTX slide created successfully.")
        return pptx_buffer

    except Exception as pptx_e:
        logger.error(f"Error during create_proper_pptx function: {pptx_e}", exc_info=True)
        return None


# --- Download Link Helper ---
def get_download_link(buffer, filename, text):
     """Generate a download link for the given file buffer"""
     try:
         if not buffer:
             logger.error("Cannot create download link: buffer is None.")
             return f'<span class="error-box">PPTX generation failed. Cannot create download link.</span>'
         if isinstance(buffer, io.BytesIO): file_bytes = buffer.getvalue()
         elif isinstance(buffer, bytes): file_bytes = buffer
         else: raise TypeError("Buffer must be bytes or BytesIO")
         b64 = base64.b64encode(file_bytes).decode(); mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
         href = f'<a href="data:{mime};base64,{b64}" download="{filename}" class="download-button">{text}</a>'
         return href
     except Exception as e:
         logger.error(f"Download link error: {e}", exc_info=True)
         return f'<span class="error-box">Error creating download link.</span>'

# --- Slide Preview Function (Restored & Adjusted) ---
def generate_slide_preview(slide_data):
    """ Generates a simplified PIL image preview of the target slide layout. """
    preview_width = 1000; preview_height = int(preview_width * 9 / 16);
    image = Image.new('RGB', (preview_width, preview_height), color=hex_to_rgb(PDQ_COLORS["revolver"]));
    draw = ImageDraw.Draw(image)
    try:
        font_bold = get_font("Segoe UI", 24, bold=True)
        font_reg = get_font("Segoe UI", 14)
        font_small = get_font("Segoe UI", 11)
        font_tiny_bold = get_font("Segoe UI", 9, bold=True)
    except Exception as font_e:
        logger.warning(f"Error loading fonts for preview, using defaults: {font_e}")
        font_bold=ImageFont.load_default(); font_reg=ImageFont.load_default();
        font_small=ImageFont.load_default(); font_tiny_bold=ImageFont.load_default()

    LEFT_MARGIN_PX=30; TOP_MARGIN_PX=30; GRID_GAP_PX=10; LOGO_SIZE_PX=40;
    TITLE_LEFT_PX=LEFT_MARGIN_PX+LOGO_SIZE_PX+15;
    RIGHT_COL_LEFT_PX=preview_width*0.58;
    RIGHT_COL_WIDTH_PX=preview_width-RIGHT_COL_LEFT_PX-LEFT_MARGIN_PX;
    BOX_BG=hex_to_rgb(PDQ_COLORS["valentino"]);
    TEXT_COLOR=hex_to_rgb(PDQ_COLORS["magnolia"]);
    TITLE_COLOR=hex_to_rgb(PDQ_COLORS["melrose"])
    BOTTOM_MARGIN = 30

    draw.rounded_rectangle([LEFT_MARGIN_PX, TOP_MARGIN_PX, LEFT_MARGIN_PX+LOGO_SIZE_PX, TOP_MARGIN_PX+LOGO_SIZE_PX], radius=5, fill=hex_to_rgb(PDQ_COLORS["electric_violet"]))
    draw.text((TITLE_LEFT_PX, TOP_MARGIN_PX + 5), slide_data.get('title', 'A/B Test'), fill=TEXT_COLOR, font=font_bold)

    hyp_top=TOP_MARGIN_PX+LOGO_SIZE_PX+30; hyp_width=RIGHT_COL_LEFT_PX-LEFT_MARGIN_PX-GRID_GAP_PX; hyp_height=100;
    draw.rounded_rectangle([LEFT_MARGIN_PX, hyp_top, LEFT_MARGIN_PX+hyp_width, hyp_top+hyp_height], radius=8, fill=BOX_BG);
    draw.text((LEFT_MARGIN_PX+15, hyp_top+10), "✓ Hypothesis", fill=TITLE_COLOR, font=font_reg)
    hyp_text=slide_data.get('hypothesis', '...')[:100] + ('...' if len(slide_data.get('hypothesis', '')) > 100 else '');
    draw.text((LEFT_MARGIN_PX+15, hyp_top+40), hyp_text, fill=TEXT_COLOR, font=font_small)

    grid_top=hyp_top+hyp_height+GRID_GAP_PX; box_w=(hyp_width-GRID_GAP_PX)/2; box_h=70;
    draw.rounded_rectangle([LEFT_MARGIN_PX, grid_top, LEFT_MARGIN_PX+box_w, grid_top+box_h], radius=8, fill=BOX_BG);
    draw.text((LEFT_MARGIN_PX+15, grid_top+10),"👥 Segment", fill=TITLE_COLOR, font=font_reg);
    draw.rounded_rectangle([LEFT_MARGIN_PX+box_w+GRID_GAP_PX, grid_top, LEFT_MARGIN_PX+hyp_width, grid_top+box_h], radius=8, fill=BOX_BG);
    draw.text((LEFT_MARGIN_PX+box_w+GRID_GAP_PX+15, grid_top+10),"📅 Timeline", fill=TITLE_COLOR, font=font_reg)

    # Adjusted heights for preview to reflect PPTX changes
    control_top=TOP_MARGIN_PX+LOGO_SIZE_PX+30;
    control_height=int((preview_height-control_top-BOTTOM_MARGIN-GRID_GAP_PX)/2 * 1.05); # Slightly taller
    variant_top=control_top+control_height+GRID_GAP_PX+30

    draw.text((RIGHT_COL_LEFT_PX, control_top-25), "Control", fill=hex_to_rgb(PDQ_COLORS["black"]), font=font_reg)
    draw.rounded_rectangle([RIGHT_COL_LEFT_PX, control_top, RIGHT_COL_LEFT_PX+RIGHT_COL_WIDTH_PX, control_top+control_height], radius=8, fill=hex_to_rgb(PDQ_COLORS["white"]), outline=hex_to_rgb(PDQ_COLORS["moon_raker"]));
    ctrl_img_preview = slide_data.get('control_image')
    if ctrl_img_preview and isinstance(ctrl_img_preview, Image.Image):
        ctrl_img_preview.thumbnail((int(RIGHT_COL_WIDTH_PX*0.8), int(control_height*0.8)))
        img_x = int(RIGHT_COL_LEFT_PX + (RIGHT_COL_WIDTH_PX - ctrl_img_preview.width) / 2)
        img_y = int(control_top + (control_height - ctrl_img_preview.height) / 2)
        image.paste(ctrl_img_preview, (img_x, img_y))
    else:
        draw.text((RIGHT_COL_LEFT_PX+20, control_top+20), "(Control Img Preview)", fill=hex_to_rgb(PDQ_COLORS["grey_text"]), font=font_reg)


    draw.text((RIGHT_COL_LEFT_PX, variant_top-25), "Variant B (example)", fill=hex_to_rgb(PDQ_COLORS["black"]), font=font_reg)
    draw.rounded_rectangle([RIGHT_COL_LEFT_PX, variant_top, RIGHT_COL_LEFT_PX+RIGHT_COL_WIDTH_PX, variant_top+control_height], radius=8, fill=hex_to_rgb(PDQ_COLORS["white"]), outline=hex_to_rgb(PDQ_COLORS["moon_raker"]));
    var_img_preview = slide_data.get('variant_image')
    if var_img_preview and isinstance(var_img_preview, Image.Image):
        var_img_preview.thumbnail((int(RIGHT_COL_WIDTH_PX*0.8), int(control_height*0.8)))
        img_x = int(RIGHT_COL_LEFT_PX + (RIGHT_COL_WIDTH_PX - var_img_preview.width) / 2)
        img_y = int(variant_top + (control_height - var_img_preview.height) / 2)
        image.paste(var_img_preview, (img_x, img_y))
    else:
        draw.text((RIGHT_COL_LEFT_PX+20, variant_top+20), "(Variant Img Preview)", fill=hex_to_rgb(PDQ_COLORS["grey_text"]), font=font_reg)

    # --- VARIANT TAG ALREADY REMOVED FROM PREVIEW ---

    draw.text((preview_width-250, preview_height-25), "PDQ A/B Test | ... | Confidential", fill=TITLE_COLOR, font=font_small)

    logger.info("Generated slide preview image.")
    return image

def create_price_variant(old_price, new_price):
    """
    Create control and variant images based on the HTML template with different prices.
    This function ensures we use the HTML-based generation for a consistent look.
    
    Args:
        old_price: Original price string (e.g. "$7.95")
        new_price: New price string (e.g. "$5.00")
        
    Returns:
        tuple: (control_image, variant_image) as PIL Image objects
    """
    logger.info(f"Creating HTML-based variant image by changing price from {old_price} to {new_price}")
    try:
        # Generate a unique timestamp to ensure we get fresh images
        timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')
        
        # Generate unique filenames with timestamp
        control_filename = f"control_shipping_{timestamp}.png"
        variant_filename = f"variant_shipping_{timestamp}.png"
        
        # Create HTML for both control and variant
        control_html = generate_shipping_html(old_price, "$24.95", is_variant=False)
        variant_html = generate_shipping_html(new_price, "$24.95", is_variant=True)
        
        logger.info("Converting HTML to images with fresh timestamp...")
        
        # Generate the images from HTML
        control_image = html_to_image(control_html, output_path=control_filename)
        variant_image = html_to_image(variant_html, output_path=variant_filename)
        
        # Log the image sizes for debugging
        logger.info(f"Generated control image with size: {control_image.size if control_image else 'None'}")
        logger.info(f"Generated variant image with size: {variant_image.size if variant_image else 'None'}")
        
        return control_image, variant_image
        
    except Exception as e:
        logger.error(f"Error in create_price_variant: {e}", exc_info=True)
        # In case of error, return placeholder images
        placeholder = Image.new('RGB', (600, 350), color=hex_to_rgb(PDQ_COLORS["moon_raker"]))
        draw = ImageDraw.Draw(placeholder)
        try:
            font = get_font("Arial", 14)
            draw.text((10, 10), f"Error generating image. Old: {old_price}, New: {new_price}", 
                     fill=(0,0,0), font=font)
        except:
            # Fallback if font fails
            draw.text((10, 10), f"Error generating image", fill=(0,0,0))
        
        control_copy = placeholder.copy()
        variant_copy = placeholder.copy()
        draw_v = ImageDraw.Draw(variant_copy)
        draw_v.text((10, 30), "VARIANT", fill=(0,0,0))
        
        return control_copy, variant_copy

# --- Streamlit UI and Main Logic ---
# Initialize session state keys if they don't exist
if 'control_image_pil' not in st.session_state:
    st.session_state.control_image_pil = None
if 'variant_image_pil' not in st.session_state: # Keep track of the variant image in state too
    st.session_state.variant_image_pil = None
if 'control_ocr_text' not in st.session_state:
    st.session_state.control_ocr_text = ""
if 'supporting_data_image_pil' not in st.session_state:
    st.session_state.supporting_data_image_pil = None
if 'supporting_data_ocr_text' not in st.session_state:
    st.session_state.supporting_data_ocr_text = ""
if 'extracted_metrics' not in st.session_state:
    st.session_state.extracted_metrics = {}
if 'pdf_text_content' not in st.session_state:
    st.session_state.pdf_text_content = []
if 'control_price' not in st.session_state:
    st.session_state.control_price = "$7.95" # Default
if 'variant_price' not in st.session_state:
    st.session_state.variant_price = "$5.00" # Default
if 'test_name' not in st.session_state:
    st.session_state.test_name = ""
if 'test_type' not in st.session_state:
    st.session_state.test_type = ""
if 'test_hypothesis' not in st.session_state:
    st.session_state.test_hypothesis = ""
if 'processed_data_file_id' not in st.session_state:
    st.session_state.processed_data_file_id = None # Track the ID of the last processed supporting data file
# Remove the 'control_image_filename' if it was just for tracking, the PIL object itself can be tracked.


st.title("🧪 PDQ A/B Test Slide Generator")
st.markdown("Generate professional A/B test slides matching the PDQ standard layout.")
st.markdown("---")

# Display general errors if any occurred during the last run
if st.session_state.get("error_message"):
    st.error(f"❌ An error occurred: {st.session_state.error_message}")

st.sidebar.header("📥 Input Parameters")
supporting_data_file = st.sidebar.file_uploader("1. Upload Supporting Data (PNG or PDF)", type=["png", "pdf"])
control_layout_file = st.sidebar.file_uploader("2. Upload Control Layout Image (PNG)", type=["png"])
if control_layout_file is not None:
    current_control_filename = getattr(control_layout_file, 'name', None)
    previous_filename = st.session_state.get('control_image_filename')
    
    if current_control_filename != previous_filename:
        # Clear all cached data when a new file is uploaded
        st.session_state.control_image_filename = current_control_filename
        
        # Clear any previously generated images
        if 'slide_data' in st.session_state:
            st.session_state.slide_data = {}
        
        # Force cache invalidation for image generation
        for key in ['shipping_control_img', 'shipping_variant_img']:
            if key in st.session_state:
                st.session_state.pop(key, None)
                
        logger.info(f"New control image detected: {current_control_filename}. Cache cleared.")

if control_layout_file:
    st.sidebar.info("Note: Your uploaded image is used for analysis only. The actual slide will use a generated shipping option mockup.")
segment = st.sidebar.text_input("3. Target Segment", placeholder="e.g., First-time mobile users")
test_type = st.sidebar.text_input("4. Test Description", placeholder="e.g., Price Test — Control: $7.95 | Variant: $5.00")
with st.sidebar.expander("🔧 Advanced Options"):
    custom_hypothesis = st.text_area("Custom Hypothesis (Optional)")
    show_debug = st.checkbox("Show Debug Information", value=False)

required_inputs_present = bool(test_type and segment and control_layout_file)
generate_button = st.sidebar.button("🚀 Generate A/B Test Slide", type="primary", disabled=not required_inputs_present, use_container_width=True)
if not required_inputs_present: st.sidebar.warning("Please provide all required inputs (2, 3, 4). Supporting data (1) is optional.")

# --- Combined Fix for both the price issue and fresh image handling ---

if generate_button:
    # CLEAR ALL PREVIOUS STATE - be thorough to avoid caching issues
    # We need to completely reset all state variables to ensure clean processing
    for key in list(st.session_state.keys()):
        if key not in ['control_image_filename']:  # Keep only the filename tracking
            st.session_state.pop(key, None)
    
    # Initialize fresh state
    st.session_state.slide_generated = False
    st.session_state.output_buffer = None
    st.session_state.slide_data = {}
    st.session_state.error_message = None
    
    logger.info("Generate button clicked. Starting process with completely fresh state...")

    with st.spinner("⚙️ Processing inputs and generating slide..."):
        try: # Wrap the entire generation process
            logger.info("Initializing SlideGeneratorHelper...")
            slide_helper = PDQSlideGeneratorHelper()
            default_metrics = { "conversion_rate": "N/A", "total_checkout": "N/A", "checkouts": "N/A", "orders": "N/A", "shipping_revenue": "N/A", "aov": "N/A" }
            metrics = default_metrics.copy()
            extracted_supporting_data_text = ""; supporting_data_image = None

            # --- Process Supporting Data (Optional) ---
            if supporting_data_file:
                logger.info(f"Processing supporting data file: {supporting_data_file.name}")
                st.sidebar.info(f"Processing '{supporting_data_file.name}'...")
                if supporting_data_file.type == "image/png":
                    extracted_supporting_data_text, img_pil = extract_text_from_image(supporting_data_file)
                    if img_pil:
                        metrics = extract_metrics_from_supporting_data(img_pil)
                        supporting_data_image = img_pil
                        logger.info("Processed PNG supporting data.")
                    else:
                        logger.warning("Failed to get PIL image from supporting PNG.")
                elif supporting_data_file.type == "application/pdf":
                    pdf_content = extract_from_pdf(supporting_data_file)
                    if pdf_content:
                        extracted_supporting_data_text = " ".join(p["text"] for p in pdf_content if p.get("text"))
                        first_image = next((img for p in pdf_content for img in p.get("images", []) if isinstance(img, Image.Image)), None)
                        if first_image:
                            supporting_data_image = first_image
                            metrics = extract_metrics_from_supporting_data(first_image)
                            logger.info("Processed PDF supporting data, found image.")
                        else:
                            logger.warning("No valid images found in the PDF for metrics/display.")
                    else:
                         logger.warning("extract_from_pdf returned empty content.")
                st.sidebar.success("Supporting data processed.")
            else:
                logger.info("No supporting data file provided.")

            # --- Process Control Layout Image and Extract Prices ---
            logger.info("Processing control layout file...")
            if not control_layout_file:
                logger.error("No control layout file provided!")
                st.error("Control Layout Image is required. Please upload an image file.")
                st.session_state["error_message"] = "Missing Control Layout Image."
                st.stop()

            # Always read fresh bytes from the uploaded file
            control_file_bytes = control_layout_file.getvalue()
            if not control_file_bytes:
                logger.error("Control layout file is empty!")
                st.error("The uploaded Control Layout Image appears to be empty.")
                st.session_state["error_message"] = "Empty Control Layout Image."
                st.stop()
                
            # Create a fresh BytesIO object from the file bytes
            control_bytes_io = io.BytesIO(control_file_bytes)

            # Open a completely fresh image from the bytes
            try:
                control_image_input_pil = Image.open(control_bytes_io)
                # Make a copy to ensure we don't have reference issues
                control_image_input_pil = control_image_input_pil.copy()
                logger.info(f"Successfully opened fresh control image: {control_image_input_pil.size}")
            except Exception as img_error:
                logger.error(f"Failed to open control image: {img_error}")
                st.error("Failed to open the Control Layout Image file.")
                st.session_state["error_message"] = f"Control Image Error: {img_error}"
                st.stop()

            # Extract text from the image directly here
            extracted_text = pytesseract.image_to_string(
                cv2.threshold(
                    cv2.cvtColor(
                        np.array(control_image_input_pil), 
                        cv2.COLOR_RGB2GRAY
                    ), 
                    0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU
                )[1]
            )
            logger.info(f"Extracted {len(extracted_text)} characters from control image")

            # Store the current filename to track changes
            current_filename = getattr(control_layout_file, 'name', 'unknown')
            previous_filename = st.session_state.get('control_image_filename')
            if current_filename != previous_filename:
                logger.info(f"New control image detected: '{current_filename}' (previously: '{previous_filename}')")
                st.session_state["control_image_filename"] = current_filename

            # --- Extract prices from test description ---
            logger.info("Extracting prices for shipping options...")
            old_price_str, new_price_str = extract_prices_from_test_type(test_type)
            logger.info(f"Extracted prices from test description: Control={old_price_str}, Variant={new_price_str}")

            # --- Create control and variant images based on the uploaded image ---
            logger.info("Creating control and variant images based on uploaded image with price changes")

            # Start with copies of the uploaded image
            control_shipping_img = control_image_input_pil.copy()
            variant_shipping_img = control_image_input_pil.copy()

            # Function to find and replace price text in an image
            
            def replace_price_in_image(img, old_price_pattern, new_price_text):
                """Find and replace price text in the image using draw operations"""
                img_copy = img.copy()
                img_width, img_height = img_copy.size
                
                # Attempt to find price locations using OCR
                # First convert to grayscale for better text detection
                img_gray = cv2.cvtColor(np.array(img_copy), cv2.COLOR_RGB2GRAY)
                
                # Use pytesseract to get detailed OCR data including bounding boxes
                ocr_data = pytesseract.image_to_data(img_gray, output_type=pytesseract.Output.DICT)
                
                # Look for price patterns in the OCR text
                price_found = False
                for i, text in enumerate(ocr_data['text']):
                    # Check if this text matches a price pattern - either exact match or contains $ symbol
                    if text and ('$' in text or re.match(r'\$?\d+\.\d+', text)):
                        logger.info(f"Found potential price text: '{text}' at index {i}")
                        
                        # Get the bounding box for this text
                        x = ocr_data['left'][i]
                        y = ocr_data['top'][i]
                        w = ocr_data['width'][i]
                        h = ocr_data['height'][i]
                        
                        # Check if the box is sensible
                        if w > 0 and h > 0 and x >= 0 and y >= 0:
                            # Create a drawable canvas
                            draw = ImageDraw.Draw(img_copy)
                            
                            # Create a background rectangle to cover the old text
                            padding = 2  # Add padding around text
                            draw.rectangle(
                                [(x-padding, y-padding), (x+w+padding, y+h+padding)],
                                fill=(255, 255, 255, 0)  # White background
                            )
                            
                            # Get a font similar to what's in the image
                            try:
                                font_size = h
                                font = ImageFont.truetype("Arial", font_size)
                            except:
                                font = ImageFont.load_default()
                            
                            # Draw the new price text
                            draw.text((x, y), new_price_text, fill=(0, 0, 0), font=font)
                            price_found = True
                            logger.info(f"Replaced price '{text}' with '{new_price_text}' at position ({x}, {y})")
                            break  # Just replace the first found price for simplicity
                
                if not price_found:
                    logger.warning(f"No price text found in the image. Adding price overlay.")
                    # If no price was found, add an overlay with the price in a visible location
                    draw = ImageDraw.Draw(img_copy)
                    try:
                        font = ImageFont.truetype("Arial", 24)
                    except:
                        font = ImageFont.load_default()
                    
                    # Add a price overlay in the top right
                    overlay_text = f"Price: {new_price_text}"
                    text_width = font.getbbox(overlay_text)[2] if hasattr(font, 'getbbox') else font.getsize(overlay_text)[0]
                    text_x = img_width - text_width - 20
                    text_y = 20
                    
                    # Add background for readability
                    padding = 5
                    draw.rectangle(
                        [(text_x-padding, text_y-padding), 
                        (text_x+text_width+padding, text_y+font.getbbox(overlay_text)[3]+padding if hasattr(font, 'getbbox') else text_y+font.getsize(overlay_text)[1]+padding)],
                        fill=(244, 250, 252),
                        outline=(0, 0, 0)
                    )
                    
                    # Draw the text
                    draw.text((text_x, text_y), overlay_text, fill=(0, 0, 0), font=font)
                
                return img_copy

            # Try to replace prices in both images
            try:
                # Use the price from the test description for the control image
                control_shipping_img = replace_price_in_image(control_shipping_img, None, old_price_str)
                logger.info(f"Set control image price to {old_price_str}")
                
                # Use the variant price for the variant image
                variant_shipping_img = replace_price_in_image(variant_shipping_img, None, new_price_str)
                logger.info(f"Set variant image price to {new_price_str}")
                
                # Add a VARIANT label to the variant image
                variant_width, variant_height = variant_shipping_img.size
                draw = ImageDraw.Draw(variant_shipping_img)
                
                # Add "VARIANT" label at the top right corner
                try:
                    font = ImageFont.truetype("Arial", 16)
                except:
                    font = ImageFont.load_default()
                    
                # Draw a rectangle with text for the variant indicator
                rect_width = 82
                rect_height = 23
                rect_x = variant_width - rect_width - 10  # 10px from the right edge
                rect_y = 10  # 10px from the top
                
                # White background
                draw.rectangle(
                    [(rect_x, rect_y), (rect_x + rect_width, rect_y + rect_height)],
                    fill=(255, 255, 255),
                    outline=(255, 0, 0),
                    width=2
                )
                
                # Add "VARIANT" text in red
                text_width = font.getbbox("VARIANT")[2] if hasattr(font, 'getbbox') else font.getsize("VARIANT")[0]
                text_x = rect_x + (rect_width - text_width) // 2
                text_y = rect_y + 5
                draw.text((text_x, text_y), "VARIANT", fill=(255, 0, 0), font=font)
                
                logger.info("Added VARIANT label to the variant image")
                image_source_method = "uploaded_with_price_edit"
                
            except Exception as e:
                logger.error(f"Error modifying uploaded images with prices: {e}")
                logger.info("Falling back to generating fresh shipping images with prices")
                
                # Fallback to generated images if price replacement fails
                control_shipping_img, variant_shipping_img = generate_shipping_options(old_price_str, new_price_str)
                image_source_method = "generated_fallback"

            logger.info(f"Successfully created control and variant images with prices: Control={old_price_str}, Variant={new_price_str}")

            # --- Generate Slide Content ---
            logger.info("Generating slide text content (hypothesis, kpis, etc.)...")
            parsed_title = slide_helper.parse_test_type(test_type)
            hypothesis = custom_hypothesis if custom_hypothesis else slide_helper.generate_hypothesis(test_type, segment, extracted_supporting_data_text)
            goal, kpi, impact = slide_helper.infer_goals_and_kpis(test_type)
            tags = slide_helper.generate_tags(test_type, segment, extracted_supporting_data_text)
            success_criteria = slide_helper.determine_success_criteria(test_type, kpi, goal)
            timeline_str = "4 weeks\nStat Sig: 85%"
            checkouts_required_str = "20,000"

            # --- Create PowerPoint ---
            logger.info("Calling create_proper_pptx function...")
            output_buffer = create_proper_pptx(
                 title=f"AB Test: {parsed_title}",
                 hypothesis=hypothesis,
                 segment=segment,
                 goal=goal,
                 kpi_impact_str=f"{kpi} ({impact} Improvement)",
                 elements_tags=tags,
                 timeline_str=timeline_str,
                 success_criteria=success_criteria,
                 checkouts_required_str=checkouts_required_str,
                 control_image=control_shipping_img,
                 variant_image=variant_shipping_img,
                 supporting_data_image=supporting_data_image
            )

            if output_buffer is None:
                logger.error("create_proper_pptx returned None. PPTX generation failed.")
                st.error("Failed to generate the PowerPoint file due to an internal error.")
                st.session_state.error_message = "PPTX generation failed (check logs for details)."
            else:
                logger.info("PPTX buffer created successfully.")
                st.session_state.slide_generated = True
                st.session_state.output_buffer = output_buffer

            # --- Update Session State (even if PPTX failed, for preview/debug) ---
            st.session_state.slide_data = {
                 "title": f"AB Test: {parsed_title}", 
                 "segment": segment, 
                 "test_type": test_type,
                 "control_image": control_shipping_img, 
                 "variant_image": variant_shipping_img,
                 "supporting_data_image": supporting_data_image, 
                 "raw_control_image": control_image_input_pil,
                 "metrics": metrics, 
                 "hypothesis": hypothesis, 
                 "goal": goal, 
                 "kpi": kpi, 
                 "impact": impact, 
                 "tags": tags, 
                 "success_criteria": success_criteria,
                 "control_price": old_price_str, 
                 "variant_price": new_price_str,
                 "image_source_method": image_source_method
            }
            logger.info(f"Session state updated with {image_source_method} images and price information.")

        except Exception as e:
            st.error(f"❌ An unexpected error occurred during slide generation: {e}")
            logger.exception("Error during slide generation button press:")
            st.session_state.slide_generated = False
            st.session_state.output_buffer = None
            st.session_state.error_message = str(e)

        logger.info("Attempting st.rerun() to update UI.")
        st.rerun()

# --- Display Results Section ---
if st.session_state.slide_generated and st.session_state.output_buffer:
    logger.info("Displaying results section (slide_generated=True, output_buffer exists)...")
    try:
        st.markdown(f'<div class="success-box">✅ A/B Test slide generated successfully!</div>', unsafe_allow_html=True)
        col1, col2 = st.columns([2, 1])

        with col1:
            st.subheader("📊 Image Previews")
            st.markdown("Previews of the images used in the generated slide.")
            logger.info("Displaying Control Image expander...")
            with st.expander("Control Image (Generated Shipping Option)", expanded=True):
                img_ctrl = st.session_state.slide_data.get('control_image')
                if img_ctrl and isinstance(img_ctrl, Image.Image):
                    st.image(img_ctrl, caption="Generated Control Shipping Image", use_column_width=True)
                    logger.info("Control image displayed.")
                else:
                    st.warning("Control image preview not available or invalid.")
                    logger.warning("Control image not available or invalid in session state for display.")

            logger.info("Displaying Variant Image expander...")
            with st.expander("Variant Image (Generated Shipping Option)", expanded=True):
                 img_var = st.session_state.slide_data.get('variant_image')
                 if img_var and isinstance(img_var, Image.Image):
                     st.image(img_var, caption="Generated Variant Shipping Image", use_column_width=True)
                     logger.info("Variant image displayed.")
                 else:
                     st.warning("Variant image preview not available or invalid.")
                     logger.warning("Variant image not available or invalid in session state for display.")

            logger.info("Displaying Supporting Data expander...")
            img_supp = st.session_state.slide_data.get('supporting_data_image')
            if img_supp and isinstance(img_supp, Image.Image):
                with st.expander("Supporting Data Image (Uploaded)", expanded=False):
                     st.image(img_supp, caption="Uploaded Supporting Data Image", use_column_width=True)
                     logger.info("Supporting data image displayed.")
                     metrics_data = st.session_state.slide_data.get('metrics', {})
                     if metrics_data and isinstance(metrics_data, dict) and any(v != "N/A" for v in metrics_data.values()):
                         st.write("**Extracted Metrics:**")
                         st.table(metrics_data)
                         logger.info("Metrics table displayed.")
            elif not supporting_data_file:
                 st.info("No supporting data image was provided.")
                 logger.info("No supporting data file provided initially.")
            else:
                 st.warning("Supporting data image preview not available (check processing logs).")
                 logger.warning("Supporting data image not available or invalid in session state for display.")

        with col2:
            st.subheader("⬇️ Download Slide")
            logger.info("Generating download link...")
            download_link_html = get_download_link(
                st.session_state.output_buffer,
                "pdq_ab_test_slide.pptx",
                "Download PPTX File"
            )
            st.markdown(download_link_html, unsafe_allow_html=True)
            logger.info("Download link displayed.")

            st.markdown("---"); st.subheader("📝 Slide Content Summary")
            logger.info("Displaying summary data...")
            summary_data = {
                "Title": st.session_state.slide_data.get("title", "N/A"),
                "Segment": st.session_state.slide_data.get("segment", "N/A"),
                "Goal": st.session_state.slide_data.get("goal", "N/A"),
                "KPI": st.session_state.slide_data.get("kpi", "N/A"),
                "Tags": ", ".join(st.session_state.slide_data.get("tags", [])),
                "Success Criteria": st.session_state.slide_data.get("success_criteria", "N/A"),
            }
            for key, value in summary_data.items():
                st.markdown(f"**{key}:** {value}")
            logger.info("Summary data displayed.")

            st.markdown("---")
            if st.button("✨ Create Another Slide"):
                logger.info("Clearing session state for new slide.")
                keys_to_clear = ['slide_generated', 'output_buffer', 'slide_data', 'error_message', 
                                'control_image_data', 'control_image_filename']
                for key in keys_to_clear:
                    st.session_state.pop(key, None)
                st.rerun()

        # --- Debug Information Display ---
        if show_debug:
            logger.info("Displaying debug information.")
            st.markdown("---"); st.subheader("🔍 Debug Information")
            debug_tabs = st.tabs(["Inputs Used", "Generated Content", "Images"])
            with debug_tabs[0]:
                st.write("Test Description Input:", st.session_state.slide_data.get('test_type', 'N/A'))
                st.write("Segment Input:", st.session_state.slide_data.get('segment', 'N/A'))
                st.write("Custom Hypothesis Input:", custom_hypothesis if custom_hypothesis else "(Not provided)")
            with debug_tabs[1]:
                st.write("Generated Hypothesis:", st.session_state.slide_data.get('hypothesis', 'N/A'))
                st.write("Inferred Goal:", st.session_state.slide_data.get('goal', 'N/A'))
                st.write("Inferred KPI:", st.session_state.slide_data.get('kpi', 'N/A'))
                st.write("Generated Tags:", st.session_state.slide_data.get('tags', []))
                st.write("Determined Success Criteria:", st.session_state.slide_data.get('success_criteria', 'N/A'))
                st.write("Extracted Metrics:", st.session_state.slide_data.get('metrics', {}))
            with debug_tabs[2]:
                 st.write("Control Image (Uploaded - Raw Input):")
                 raw_ctrl_img = st.session_state.slide_data.get('raw_control_image')
                 if raw_ctrl_img and isinstance(raw_ctrl_img, Image.Image): st.image(raw_ctrl_img, width=300)
                 else: st.write("(Not available or invalid)")

                 st.write("Supporting Data Image (Used in PPTX):")
                 supp_img = st.session_state.slide_data.get('supporting_data_image')
                 if supp_img and isinstance(supp_img, Image.Image): st.image(supp_img, width=300)
                 else: st.write("(Not available or invalid)")

                 st.write("Generated Control Shipping Image (Used in PPTX):")
                 gen_ctrl_img = st.session_state.slide_data.get('control_image')
                 if gen_ctrl_img and isinstance(gen_ctrl_img, Image.Image): st.image(gen_ctrl_img, width=300)
                 else: st.write("(Not available or invalid)")

                 st.write("Generated Variant Shipping Image (Used in PPTX):")
                 gen_var_img = st.session_state.slide_data.get('variant_image')
                 if gen_var_img and isinstance(gen_var_img, Image.Image): st.image(gen_var_img, width=300)
                 else: st.write("(Not available or invalid)")

    except Exception as display_e:
        st.error(f"❌ An error occurred while displaying the results: {display_e}")
        logger.exception("Error during results display section:")
        st.session_state.error_message = f"Error displaying results: {display_e}"
        st.session_state.slide_generated = False

# --- Fallback message if not generated ---
elif not st.session_state.error_message and not generate_button:
    # Only show the initial info message if no error is pending and the button wasn't just clicked
    st.info("⬆️ Upload files and fill in details in the sidebar to generate the slide.")
    # You can optionally add back the structure guide here if desired
    # st.markdown("##### Target Slide Structure Guide:")
    # st.markdown("...")


# --- Custom Footer ---
footer_year = datetime.datetime.now().year; footer_left_text = "PDQ A/B Test Slide Generator | Streamlining Test Documentation"; footer_right_text = f"PDQ © {footer_year}"
st.markdown(f"""<div class="custom-footer"><div class="footer-left">{footer_left_text}</div><div class="footer-right">{footer_right_text}</div></div>""", unsafe_allow_html=True)
              
              
              
